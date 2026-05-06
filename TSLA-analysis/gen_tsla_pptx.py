"""gen_tsla_pptx.py — Tesla deep valuation summary deck (11 slides)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime, os

OUT = "/root/PythonProgramming/TSLA-analysis"

# Colors
DARK  = RGBColor(0x1D, 0x1D, 0x1F)
TR    = RGBColor(0xE3, 0x19, 0x37)  # Tesla red
GRN   = RGBColor(0x34, 0xC7, 0x59)
BLU   = RGBColor(0x00, 0x71, 0xE3)
LGRY  = RGBColor(0xF5, 0xF5, 0xF7)
MGRY  = RGBColor(0x86, 0x86, 0x8B)
WHT   = RGBColor(0xFF, 0xFF, 0xFF)
AMB   = RGBColor(0xFF, 0x95, 0x00)
DK2   = RGBColor(0x2C, 0x2C, 0x2E)
RED   = RGBColor(0xD7, 0x00, 0x15)
PURP  = RGBColor(0xA8, 0x55, 0xF7)
CYAN  = RGBColor(0x00, 0xC2, 0xE0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def box(slide, text, l, t, w, h, sz=14, bold=False, fc=DARK, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = fc
    return tx

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def card(slide, l, t, w, h, title, body, title_color=DARK, accent=BLU, bg_color=WHT):
    rect(slide, l, t, w, h, bg_color)
    rect(slide, l, t, 0.18, h, accent)
    box(slide, title, l+0.25, t+0.08, w-0.35, 0.38, sz=12, bold=True, fc=title_color)
    box(slide, body,  l+0.25, t+0.52, w-0.35, h-0.62, sz=10, fc=MGRY, italic=False)


# ── S1: Title ─────────────────────────────────────────────────────────────────
s1 = blank(); bg(s1, DARK)
rect(s1, 0, 5.8, 13.33, 0.1, TR)
box(s1, "TESLA, INC. (TSLA)", 0.8, 0.8, 11.5, 1.4, sz=52, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
box(s1, "Deep Valuation Analysis", 0.8, 2.4, 11.5, 0.9, sz=28, fc=MGRY, align=PP_ALIGN.CENTER)
box(s1, f"Prepared {datetime.date.today().strftime('%B %Y')}  ·  NASDAQ: TSLA  ·  "
        "Consumer Discretionary / AI & Autonomy",
    0.8, 3.6, 11.5, 0.5, sz=14, fc=MGRY, align=PP_ALIGN.CENTER)
box(s1, "Models included: 3-Statement · DCF · Comp Analysis · Sum of the Parts · Scenario Bridge",
    0.8, 4.4, 11.5, 0.5, sz=13, fc=RGBColor(0x60, 0x60, 0x68), align=PP_ALIGN.CENTER, italic=True)


# ── S2: Investment Summary ────────────────────────────────────────────────────
s2 = blank(); bg(s2, LGRY)
rect(s2, 0, 0, 13.33, 1.1, DARK)
box(s2, "INVESTMENT SUMMARY", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

stats = [
    ("$887B", "Market Cap"),
    ("$97.7B", "FY2024 Revenue"),
    ("$12.5B", "FY2024 EBITDA"),
    ("71x", "EV / EBITDA"),
    ("~$278", "Current Price"),
    ("1.79M", "FY2024 Deliveries"),
]
for i, (val, lbl) in enumerate(stats):
    x = 0.3 + i * 2.17
    rect(s2, x, 1.25, 1.95, 1.6, WHT)
    box(s2, val, x, 1.35, 1.95, 0.75, sz=22, bold=True, fc=TR, align=PP_ALIGN.CENTER)
    box(s2, lbl, x, 2.1,  1.95, 0.6, sz=11, fc=MGRY, align=PP_ALIGN.CENTER)

thesis = [
    ("Bull Case (~$439/share)", GRN,
     "FSD reaches 10M+ subscribers. Cybercab launches at scale. "
     "Optimus ships 500K+ units. Energy becomes #1 grid-storage provider globally. "
     "Tesla transitions from auto OEM to AI/robotics platform."),
    ("Base Case (~$164/share)", BLU,
     "Energy doubles by FY2027. FSD reaches 5M subscribers at $99/mo. "
     "Cybercab pilots in 2–3 cities. Affordable model drives deliveries to 2.8M. "
     "Margins recover to 10–12% EBIT by FY2026."),
    ("Bear Case (~$64/share)", RED,
     "Chinese competition (BYD, Huawei) erodes margin. FSD regulatory hurdles delay monetization. "
     "Optimus R&D drag exceeds $5B/yr with no revenue. "
     "Auto gross margin falls to 14–15%."),
]
for i, (title, color, body) in enumerate(thesis):
    x = 0.3 + i * 4.35
    rect(s2, x, 3.05, 4.1, 3.75, DK2)
    rect(s2, x, 3.05, 0.18, 3.75, color)
    box(s2, title, x+0.28, 3.12, 3.75, 0.42, sz=13, bold=True, fc=color)
    box(s2, body,  x+0.28, 3.62, 3.75, 3.0,  sz=11, fc=WHT)


# ── S3: Business Overview ─────────────────────────────────────────────────────
s3 = blank(); bg(s3, WHT)
rect(s3, 0, 0, 13.33, 1.1, DARK)
box(s3, "BUSINESS OVERVIEW — SEGMENTS & PRODUCTS", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

segs = [
    ("Automotive\n52% → 79%* Revenue", TR,
     "• 6 models: Model 3/Y (volume), S/X (premium), Cybertruck, Semi\n"
     "• FY2024: 1.79M deliveries, $77.1B revenue\n"
     "• Declining ASP (~$43K) reflects mix shift to lower-trim models\n"
     "• Gross margin 17.9% — compressed from 28% peak in 2022\n"
     "• FSD: $8K one-time or $99/mo subscription; not yet L4"),
    ("Energy Gen & Storage\n10% Revenue", RGBColor(0x00,0xB3,0x86),
     "• Megapack: utility-scale battery (3 MWh each)\n"
     "• FY2024: $10.1B revenue, +165% in 2yr\n"
     "• Gross margin 24.6% in Q4'24 — best-in-class\n"
     "• AI data centers driving Megapack demand\n"
     "• Powerwall: residential storage; 50GWh deployed cumulative"),
    ("Services & Other\n11% Revenue", PURP,
     "• Service centers, body shops, Tesla Insurance\n"
     "• Supercharger: 65K+ connectors; licensed to Ford, GM, Rivian, BMW\n"
     "• FSD subscription revenue (nascent)\n"
     "• Gross margin 7.3% — services still a cost center\n"
     "• Path to 15-20% GM as high-margin FSD scales"),
    ("Emerging Platforms\n(Option Value)", AMB,
     "• Optimus Gen 2: humanoid robot, target $20-30K/unit\n"
     "• Cybercab: driverless robotaxi, announced for 2026\n"
     "• Dojo supercomputer: Tesla-built AI training chip\n"
     "• FSD v13 → targeting Level 4 unsupervised by 2026\n"
     "• $30T total addressable market for physical AI"),
]
for i, (title, color, body) in enumerate(segs):
    x = 0.25 + i * 3.22
    rect(s3, x, 1.2, 3.0, 5.65, LGRY)
    rect(s3, x, 1.2, 3.0, 0.55, color)
    box(s3, title, x+0.1, 1.22, 2.8, 0.5, sz=12, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
    box(s3, body, x+0.1, 1.85, 2.8, 4.9, sz=10, fc=DARK)


# ── S4: Financial Performance ─────────────────────────────────────────────────
s4 = blank(); bg(s4, LGRY)
rect(s4, 0, 0, 13.33, 1.1, DARK)
box(s4, "FINANCIAL PERFORMANCE — FY2021A–FY2029E ($M)", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

table_data = [
    ["Metric",        "FY2021A","FY2022A","FY2023A","FY2024A","FY2025E","FY2026E","FY2027E"],
    ["Revenue",       "$53,823","$81,462","$96,773","$97,690","$105,500","$130,500","$166,000"],
    ["Rev. Growth",   "+73%",   "+51%",   "+19%",   "+0.9%",  "+7.9%",  "+23.7%", "+27.2%"],
    ["Gross Profit",  "$13,606","$20,853","$17,660","$17,522","$19,500","$26,100","$36,520"],
    ["Gross Margin",  "25.3%",  "25.6%",  "18.2%",  "17.9%",  "18.5%",  "20.0%",  "22.0%"],
    ["EBIT",          " $6,523","$13,656"," $8,891"," $7,071"," $8,968","$13,050","$19,920"],
    ["EBIT Margin",   "12.1%",  "16.8%",  "9.2%",   "7.2%",   "8.5%",   "10.0%",  "12.0%"],
    ["EBITDA",        "$9,023","$17,256","$13,491","$12,471","$14,968","$20,050","$28,420"],
    ["Net Income",    " $5,519","$12,556","$15,000"," $7,092"," $8,000","$11,000","$16,800"],
    ["EPS (adj.)",    " $1.72", " $3.90", " $2.84", " $2.22", " $2.51", " $3.46", " $5.28"],
    ["Free Cash Flow"," $4,982"," $7,327"," $4,358"," $3,620"," $6,000"," $9,500","$14,000"],
    ["FCF Margin",    "9.3%",   "9.0%",   "4.5%",   "3.7%",   "5.7%",   "7.3%",   "8.4%"],
]
col_ws = [2.2, 1.55, 1.55, 1.55, 1.55, 1.55, 1.55, 1.55]
col_xs = [0.25]
for w in col_ws[:-1]: col_xs.append(col_xs[-1] + w)

for i, row in enumerate(table_data):
    y = 1.2 + i * 0.52
    is_hdr  = (i == 0)
    is_bold = row[0] in ("Revenue","Gross Profit","EBIT","EBITDA","Free Cash Flow")
    bg_c = DARK if is_hdr else (TR if is_bold else (WHT if i%2==0 else LGRY))
    for j, (cell_txt, cx, cw_) in enumerate(zip(row, col_xs, col_ws)):
        fc_ = WHT if (is_hdr or is_bold) else (TR if j == 0 else DARK)
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        rect(s4, cx, y, cw_, 0.48, bg_c)
        box(s4, cell_txt, cx+0.04, y+0.04, cw_-0.08, 0.4,
            sz=10, bold=(is_hdr or is_bold), fc=fc_, align=al)

box(s4, "† FY2023 reported NI = $15B incl. $5.9B non-recurring DTA benefit; adj. NI = $9.1B  |  "
        "FY2025–2027E: FactSet consensus / proprietary estimates",
    0.25, 7.1, 12.8, 0.35, sz=9, fc=MGRY, italic=True)


# ── S5: DCF Analysis ──────────────────────────────────────────────────────────
s5 = blank(); bg(s5, WHT)
rect(s5, 0, 0, 13.33, 1.1, DARK)
box(s5, "DCF VALUATION — BASE CASE ANALYSIS", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

# Left: assumptions
rect(s5, 0.25, 1.2, 4.2, 5.7, LGRY)
box(s5, "KEY ASSUMPTIONS", 0.4, 1.25, 4.0, 0.45, sz=13, bold=True, fc=TR)
asmp_items = [
    ("WACC",                 "10.0%"),
    ("Risk-Free Rate",       "4.35%"),
    ("Beta (levered)",       "2.10x"),
    ("Cost of Equity",       "14.9%"),
    ("Terminal Growth Rate", "4.0%"),
    ("Terminal EV/EBITDA",   "20.0x"),
    ("Tax Rate",             "10.0%"),
    ("Forecast Horizon",     "5 years"),
    ("Net Cash",             "$28.9B"),
    ("Diluted Shares",       "3,194M"),
]
for k, (lbl, val) in enumerate(asmp_items):
    y = 1.78 + k * 0.47
    box(s5, lbl, 0.4, y, 2.5, 0.42, sz=11, fc=DARK)
    box(s5, val, 2.9, y, 1.4, 0.42, sz=11, bold=True, fc=BLU, align=PP_ALIGN.RIGHT)

# Center: UFCF table
rect(s5, 4.65, 1.2, 8.4, 5.7, DK2)
box(s5, "UNLEVERED FCF PROJECTION  &  VALUATION OUTPUT", 4.75, 1.25, 8.2, 0.4,
    sz=12, bold=True, fc=WHT)

dcf_table = [
    ["Year",     "FY25E", "FY26E", "FY27E", "FY28E", "FY29E"],
    ["Revenue",  "$106B", "$131B", "$166B", "$210B", "$260B"],
    ["EBIT",     "$9.0B", "$13.1B","$19.9B","$27.3B","$35.4B"],
    ["NOPAT",    "$8.1B", "$11.7B","$17.9B","$24.6B","$31.8B"],
    ["UFCF",     "$5.6B", "$9.2B", "$14.4B","$21.1B","$30.5B"],
    ["PV(UFCF)", "$5.1B", "$7.6B", "$10.8B","$14.4B","$18.9B"],
]
for i, row in enumerate(dcf_table):
    y = 1.75 + i * 0.58
    is_hdr = (i == 0)
    for j, txt in enumerate(row):
        x = 4.75 + j * 1.38
        fc_ = AMB if (is_hdr and j==0) else (WHT if is_hdr else (AMB if j==0 else WHT))
        box(s5, txt, x, y, 1.3, 0.52, sz=11, bold=is_hdr,
            fc=fc_, align=PP_ALIGN.RIGHT if j>0 else PP_ALIGN.LEFT)

box(s5, "Sum of PV(UFCF) = $56.8B", 4.75, 5.35, 8.2, 0.4, sz=12, fc=AMB, bold=True)

output_rows = [
    ("Terminal Value (20x exit EV/EBITDA)", "$965B → PV $599B"),
    ("Enterprise Value (exit multiple)",    "$656B"),
    ("(+) Net Cash",                        "$28.9B"),
    ("Equity Value",                        "$685B"),
    ("Implied Share Price",                 "~$215/share"),
    ("GGM Terminal (g=4%) — cross-check",   "~$195/share"),
]
for k, (lbl, val) in enumerate(output_rows):
    y = 5.85 + k * 0.0
    is_px = "Implied" in lbl
    pass

box(s5, "Valuation Output:  Base EV $656B  →  Equity $685B  →  $215/share  "
        "(GGM cross-check: $195/share)",
    4.75, 5.78, 8.2, 0.42, sz=12, bold=True, fc=TR)
box(s5, "Current price $278 implies WACC ~8.5% or terminal multiple of 26x — "
        "market is pricing significant optionality beyond 5-yr horizon.",
    4.75, 6.28, 8.2, 0.7, sz=11, fc=MGRY, italic=True)


# ── S6: SOTP ──────────────────────────────────────────────────────────────────
s6 = blank(); bg(s6, LGRY)
rect(s6, 0, 0, 13.33, 1.1, DARK)
box(s6, "SUM OF THE PARTS — SEGMENT VALUATION", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

sotp_segs = [
    ("Automotive\n(ex-FSD)", TR,      "$97B",  "$139B", "$187B", "8x FY2027E EBITDA"),
    ("Energy Storage",       RGBColor(0x00,0xB3,0x86), "$65B", "$116B", "$175B", "22x FY2027E EBITDA"),
    ("FSD / Software",       PURP,    "$0",    "$100B", "$300B", "25-40x FY2027E Revenue"),
    ("Cybercab / Robotaxi",  AMB,     "$0",    "$30B",  "$150B", "Network optionality"),
    ("Supercharger",         CYAN,    "$12B",  "$24B",  "$40B",  "20x FY2027E EBITDA"),
    ("Insurance & Other",    BLU,     "$3B",   "$10B",  "$20B",  "10x revenue"),
    ("Optimus Robots",       RGBColor(0xFF,0x45,0x22), "$0", "$75B", "$500B", "DCF / option value"),
]

box(s6, "Segment",        0.25, 1.2, 2.5,  0.42, sz=11, bold=True, fc=DARK)
box(s6, "Bear EV",        2.85, 1.2, 1.55, 0.42, sz=11, bold=True, fc=RED,  align=PP_ALIGN.CENTER)
box(s6, "Base EV",        4.5,  1.2, 1.55, 0.42, sz=11, bold=True, fc=BLU,  align=PP_ALIGN.CENTER)
box(s6, "Bull EV",        6.15, 1.2, 1.55, 0.42, sz=11, bold=True, fc=GRN,  align=PP_ALIGN.CENTER)
box(s6, "Valuation Basis",7.8,  1.2, 5.0,  0.42, sz=11, bold=True, fc=DARK)

for i, (seg, color, bear, base, bull, basis) in enumerate(sotp_segs):
    y = 1.7 + i * 0.72
    bg_c = WHT if i % 2 == 0 else LGRY
    rect(s6, 0.25, y, 12.8, 0.68, bg_c)
    rect(s6, 0.25, y, 0.12, 0.68, color)
    box(s6, seg,   0.45, y+0.05, 2.4,  0.58, sz=11, bold=True, fc=DARK)
    box(s6, bear,  2.85, y+0.12, 1.55, 0.42, sz=11, fc=RED,  align=PP_ALIGN.CENTER)
    box(s6, base,  4.5,  y+0.12, 1.55, 0.42, sz=11, bold=True, fc=DARK, align=PP_ALIGN.CENTER)
    box(s6, bull,  6.15, y+0.12, 1.55, 0.42, sz=11, fc=GRN,  align=PP_ALIGN.CENTER)
    box(s6, basis, 7.8,  y+0.1,  4.8,  0.5,  sz=10, fc=MGRY, italic=True)

# Totals
y = 7.0
rect(s6, 0.25, y, 12.8, 0.42, DARK)
box(s6, "EQUITY VALUE (incl. $28.9B net cash)", 0.4, y+0.05, 2.4, 0.32, sz=11, bold=True, fc=WHT)
for x, val, fc_ in [(2.85,"$235B",RED),(4.5,"$551B",AMB),(6.15,"$1,401B",GRN)]:
    box(s6, val, x, y+0.05, 1.55, 0.32, sz=12, bold=True, fc=fc_, align=PP_ALIGN.CENTER)
for x, val, fc_ in [(2.85,"~$74",RED),(4.5,"~$173",AMB),(6.15,"~$439",GRN)]:
    box(s6, f"→ ${val.replace('~$','')}/sh", x, y+0.3, 1.55, 0.32, sz=10, bold=True, fc=fc_, align=PP_ALIGN.CENTER)


# ── S7: Comps ─────────────────────────────────────────────────────────────────
s7 = blank(); bg(s7, WHT)
rect(s7, 0, 0, 13.33, 1.1, DARK)
box(s7, "COMPARABLE COMPANY ANALYSIS", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

comp_data = [
    ["Company",    "Mkt Cap","EV/Rev","EV/EBITDA","P/E NTM","Rev Growth","EBIT Margin","Commentary"],
    ["Tesla","$887B","8.8x","68.8x","125x","+0.9%","7.2%","Priced as AI platform"],
    ["BYD",       "$88B",  "0.9x", "8.5x",  "22x", "+15%","6.0%","China EV leader"],
    ["GM",        "$56B",  "0.5x", "3.1x",  "5.2x","+0.9%","6.2%","Legacy scale"],
    ["Ford",      "$45B",  "0.5x", "4.0x",  "9.1x","+4.8%","4.0%","EV transition"],
    ["BMW",       "$55B",  "0.4x", "3.5x",  "5.1x","-1.2%","7.3%","Premium ICE"],
    ["NVIDIA",    "$3.2T","25.0x","48.5x",  "38x","+78%", "62.0%","AI hardware leader"],
    ["Alphabet",  "$2.0T", "5.5x","17.9x",  "21x","+14%", "31.5%","Waymo competitor"],
    ["Mobileye",  "$12B",  "4.5x","25.0x",  "45x","+14%", "18.0%","ADAS direct comp"],
    ["Auto Median","—",    "0.5x", "3.8x",  "5.6x", "3.0%","6.1%",""],
    ["Tech Median","—",    "5.5x","25.0x",  "41x","+46%", "40.0%",""],
    ["TSLA Premium to Auto","—","17×","18×","22×","—","—","13–68x premium"],
]
cw2 = [2.1, 1.3, 1.2, 1.3, 1.2, 1.3, 1.3, 3.5]
cx2 = [0.25]
for w in cw2[:-1]: cx2.append(cx2[-1]+w)

for i, row in enumerate(comp_data):
    y2 = 1.2 + i * 0.55
    is_hdr   = (i == 0)
    is_tsla  = (i == 1)
    is_stat  = (row[0] in ("Auto Median","Tech Median"))
    is_prem  = "Premium" in row[0]
    bg_c = DARK if is_hdr else (TR if is_tsla else (DK2 if is_stat else (AMB if is_prem else (LGRY if i%2==0 else WHT))))
    fc_l = WHT if (is_hdr or is_tsla or is_stat or is_prem) else DARK
    for j, (txt, cx_, cw_) in enumerate(zip(row, cx2, cw2)):
        rect(s7, cx_, y2, cw_, 0.52, bg_c)
        al = PP_ALIGN.LEFT if j in (0, 7) else PP_ALIGN.RIGHT
        fc_ = TR if (is_stat and j>0) else (GRN if is_prem else fc_l)
        box(s7, txt, cx_+0.04, y2+0.06, cw_-0.08, 0.42, sz=10,
            bold=(is_hdr or is_tsla or is_stat or is_prem), fc=fc_, align=al)

box(s7, "Tesla trades at 8.8x EV/Revenue and 69x EV/EBITDA — aligned with AI/software peers, not auto. "
        "Market assigns ~$700B of $887B market cap to non-auto optionality.",
    0.25, 7.15, 12.8, 0.35, sz=10, fc=MGRY, italic=True)


# ── S8: Bear / Base / Bull ────────────────────────────────────────────────────
s8 = blank(); bg(s8, DARK)
rect(s8, 0, 0, 13.33, 1.1, TR)
box(s8, "SCENARIO ANALYSIS — BEAR · BASE · BULL", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

scenarios = [
    ("BEAR CASE\n~$64/share", RED,
     "FY2029 Revenue: $175B\nEBIT Margin: 9.0%\nEBITDA: $26B\n\n"
     "Assumptions:\n"
     "• China competition erodes auto margins to 15%\n"
     "• FSD stuck at L2+; no monetization breakthrough\n"
     "• Optimus delayed; R&D drag $3-5B/yr\n"
     "• Energy growth slows to 10%/yr\n"
     "• Deliveries 2.2M (affordable model delayed)\n\n"
     "SOTP: Auto $97B + Energy $65B + no optionality\n"
     "= EV $174B + cash = Equity $203B → $64/sh"),
    ("BASE CASE\n~$164/share", BLU,
     "FY2029 Revenue: $260B\nEBIT Margin: 13.6%\nEBITDA: $48B\n\n"
     "Assumptions:\n"
     "• Affordable model launches FY2026; 2.8M deliveries\n"
     "• FSD subscription: 5M users @ $99/mo by FY2027\n"
     "• Cybercab pilots Austin, SF, 1 international city\n"
     "• Megapack grows to $22B by FY2028 (28% margin)\n"
     "• Optimus: ~10K units initial deployment FY2027\n\n"
     "SOTP: Segments $493B + cash = Equity $522B → $164/sh"),
    ("BULL CASE\n~$439/share", GRN,
     "FY2029 Revenue: $350B+\nEBIT Margin: 18%+\nEBITDA: $80B+\n\n"
     "Assumptions:\n"
     "• FSD L4 approved nationally; OEM licensing begins\n"
     "• Robotaxi: 500K+ Cybercabs; $2B+ annual revenue\n"
     "• Optimus: 500K+ units @ $25K ASP; $12.5B revenue\n"
     "• Dojo accelerates AI training cost by 10x\n"
     "• Energy 25GWh/yr global, 30%+ margins\n\n"
     "SOTP: Segments $1.37T + cash = Equity $1.4T → $439/sh"),
]
for i, (title, color, body) in enumerate(scenarios):
    x = 0.25 + i * 4.35
    rect(s8, x, 1.25, 4.1, 5.95, DK2)
    rect(s8, x, 1.25, 4.1, 0.9, color)
    box(s8, title, x+0.15, 1.28, 3.8, 0.82, sz=15, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
    box(s8, body,  x+0.15, 2.25, 3.8, 4.85, sz=10, fc=WHT)


# ── S9: Catalysts ─────────────────────────────────────────────────────────────
s9 = blank(); bg(s9, LGRY)
rect(s9, 0, 0, 13.33, 1.1, DARK)
box(s9, "KEY CATALYSTS — POSITIVE", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

cats = [
    ("Affordable Model\n(~$25K)", GRN,
     "Expected H1 2026 launch. Target >1M incremental annual units. "
     "Largest addressable market (~$25K price point = 60% of global new car market). "
     "Could drive 50%+ revenue growth in FY2026–2027."),
    ("FSD to Level 4", RGBColor(0x34,0xAA,0xFF),
     "FSD v12/13 showing rapid improvement. Unsupervised driving in Austin pilot (2025). "
     "L4 approval = $1,000+/mo subscription or one-time $15K. At 10M subs: $10B+ annual recurring revenue."),
    ("Cybercab Robotaxi", AMB,
     "Announced June 2023, expected FY2026. No pedals/steering wheel → lower COGS. "
     "Network effects: more rides = better AI = lower accident rate. "
     "Rides: ~$0.40-$0.50/mile vs. Uber ~$2.50 → massive structural cost advantage."),
    ("Optimus Robots", TR,
     "Gen 2 (2024): 30 kph, 20 kg payload. Target $20–30K per unit. "
     "$30T robot labor TAM. Tesla's advantages: vertical integration, AI training data, "
     "manufacturing scale (Model Y tooling learnings). Elon target: 1B units long-term."),
    ("Megapack / Grid Scale", RGBColor(0x00,0xB3,0x86),
     "AI data centers require 24/7 power + storage. Megapack backlog >$10B. "
     "IRA incentives extend 10yr. Utility procurement contracts signed in AU, EU, US. "
     "24.6% gross margin in Q4'24 — expanding as Giga Nevada 2 ramps."),
    ("Dojo + AI Platform", PURP,
     "Custom AI training chip (D1) designed for FSD training. "
     "Could offer AI-as-a-service to third parties. "
     "Real-world video data from 5B+ miles of fleet driving = "
     "most valuable autonomous driving dataset in the world."),
]
for i, (title, color, body) in enumerate(cats):
    col = i % 2; row = i // 2
    x = 0.25 + col * 6.55; y = 1.2 + row * 2.1
    rect(s9, x, y, 6.2, 1.95, WHT)
    rect(s9, x, y, 0.18, 1.95, color)
    box(s9, title, x+0.28, y+0.08, 5.8, 0.5, sz=12, bold=True, fc=DARK)
    box(s9, body,  x+0.28, y+0.62, 5.8, 1.25, sz=10, fc=MGRY)


# ── S10: Risks ────────────────────────────────────────────────────────────────
s10 = blank(); bg(s10, WHT)
rect(s10, 0, 0, 13.33, 1.1, DARK)
box(s10, "KEY RISKS", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

risks = [
    ("Chinese Competition", RED,
     "BYD sold 1.76M EVs in FY2024 vs Tesla's 1.79M. BYD Seal vs Model 3: 40% cheaper. "
     "Huawei's AITO brand growing rapidly. China = 19% of Tesla revenue at risk."),
    ("FSD Regulatory Risk", RED,
     "NHTSA investigations ongoing. FSD still classified L2 (driver required). "
     "EU has stricter type-approval rules for L3+. "
     "One fatal supervised-FSD accident could delay nationwide L4 approval by 2–5 years."),
    ("Margin Compression", RED,
     "Price cuts since 2022: Model 3 from $47K to $38K. "
     "Every 1% gross margin decline = ~$1B EBITDA impact. "
     "Structural pressure from commodity costs, labor, and competition."),
    ("Elon Musk Key-Man Risk", AMB,
     "Musk also runs SpaceX, X, xAI, Boring Co., DOGE advisory. "
     "Tesla requires his technical vision and recruiting brand. "
     "Political controversy affecting brand perception (EU, CA boycotts). "
     "Compensation package ($56B) still under legal review."),
    ("Execution on New Products", AMB,
     "Cybertruck ramp slower than expected. Semi in very limited production. "
     "Affordable model delayed 2+ years from original timeline. "
     "Optimus and Cybercab both carry high execution risk."),
    ("Valuation Multiple Risk", AMB,
     "At 69x EV/EBITDA, any guidance miss or delay re-rates the stock 30–40%. "
     "Rising interest rates compress growth multiples. "
     "If FSD or Optimus fail to scale, stock could de-rate to 20x like a traditional OEM."),
]
for i, (title, color, body) in enumerate(risks):
    col = i % 2; row = i // 2
    x = 0.25 + col * 6.55; y = 1.2 + row * 2.08
    rect(s10, x, y, 6.2, 1.95, LGRY)
    rect(s10, x, y, 0.18, 1.95, color)
    box(s10, title, x+0.28, y+0.08, 5.8, 0.5, sz=12, bold=True, fc=DARK)
    box(s10, body,  x+0.28, y+0.62, 5.8, 1.25, sz=10, fc=MGRY)


# ── S11: Valuation Summary ────────────────────────────────────────────────────
s11 = blank(); bg(s11, DARK)
rect(s11, 0, 0, 13.33, 1.1, TR)
box(s11, "VALUATION SUMMARY — PRICE TARGET RANGE", 0.3, 0.2, 12, 0.75, sz=22, bold=True, fc=WHT)

rect(s11, 0.25, 1.2, 12.8, 1.3, DK2)
box(s11, "CURRENT PRICE: ~$278/share  |  MARKET CAP: ~$887B  |  EV: ~$858B",
    0.4, 1.3, 12.5, 0.45, sz=14, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
box(s11, "At $278, market implies: Auto at fair value + ~$600B premium for FSD / Optimus / Robotaxi / Energy optionality",
    0.4, 1.75, 12.5, 0.4, sz=12, fc=AMB, italic=True, align=PP_ALIGN.CENTER)

methods = [
    ("DCF (10% WACC)", "$195–$215", "5-yr UFCF, exit multiple 20x / GGM 4%"),
    ("SOTP Bear",      "~$64",      "Auto + Energy only; no FSD/Optimus"),
    ("SOTP Base",      "~$164",     "All segments, base-case optionality"),
    ("SOTP Bull",      "~$439",     "FSD licensed + Optimus 500K+ + Robotaxi"),
    ("Auto Comps",     "$30–$80",   "3–8x FY2027E EBITDA, no tech premium"),
    ("Tech Comps",     "$300–$500", "20–30x EV/Revenue applied to $260B FY2029E"),
    ("Consensus PT",   "~$280",     "FactSet: 35 buys, 15 holds, 10 sells"),
]
rect(s11, 0.25, 2.7, 12.8, 0.52, DK2)
box(s11, "Method", 0.4, 2.77, 3.0, 0.38, sz=11, bold=True, fc=WHT)
box(s11, "Price Target", 3.5, 2.77, 2.5, 0.38, sz=11, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
box(s11, "Key Assumption", 6.1, 2.77, 6.8, 0.38, sz=11, bold=True, fc=WHT)

for i, (method, pt, note) in enumerate(methods):
    y = 3.28 + i * 0.54
    bg_c = RGBColor(0x28, 0x28, 0x2A) if i % 2 == 0 else DK2
    rect(s11, 0.25, y, 12.8, 0.5, bg_c)
    box(s11, method, 0.4, y+0.06, 3.0, 0.38, sz=11, bold=True, fc=WHT)
    fc_ = RED if "Bear" in method else (GRN if "Bull" in method else (AMB if "Consensus" in method else BLU))
    box(s11, pt, 3.5, y+0.06, 2.5, 0.38, sz=12, bold=True, fc=fc_, align=PP_ALIGN.CENTER)
    box(s11, note, 6.1, y+0.06, 6.8, 0.38, sz=10, fc=MGRY)

box(s11, "Conclusion: Tesla is valued as a call option on autonomous AI + robotics, not as an auto OEM. "
         "Fundamental downside ~$64 (bear); upside ~$440+ (bull). "
         "For long-term investors: disciplined entry below $200 offers favorable risk/reward.",
    0.25, 7.07, 12.8, 0.4, sz=11, bold=True, fc=TR, italic=True)


out = os.path.join(OUT, "TSLA_Deep_Valuation.pptx")
prs.save(out)
print(f"✓ TSLA_Deep_Valuation.pptx")
print("\n✅ All files generated.")
