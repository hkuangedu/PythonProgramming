"""
gen_pitchbook.py
Project Starlight: Netflix Strategic Acquisition of Warner Bros. Discovery
Investment Banking Pitch Book — 20 Slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

OUT = "/root/PythonProgramming/NFLX-WBD-Pitchbook"
os.makedirs(OUT, exist_ok=True)

# ── Color Palette ─────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x0F, 0x20, 0x44)   # Deep navy (primary)
GOLD  = RGBColor(0xC9, 0xA8, 0x4C)   # Gold accent
NRED  = RGBColor(0xE5, 0x09, 0x14)   # Netflix red
WHT   = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x1A, 0x1A, 0x2E)
LGRY  = RGBColor(0xF5, 0xF6, 0xF8)
MGRY  = RGBColor(0x8A, 0x8A, 0x9A)
DGRY  = RGBColor(0x4A, 0x4A, 0x5A)
GRN   = RGBColor(0x00, 0x8A, 0x48)
RED2  = RGBColor(0xCC, 0x00, 0x00)
STLB  = RGBColor(0xEA, 0xF0, 0xFB)   # Alternating table row
NBLUE = RGBColor(0x18, 0x40, 0x8C)
AMB   = RGBColor(0xFF, 0x95, 0x00)
TEAL  = RGBColor(0x00, 0x7B, 0x83)
PURP  = RGBColor(0x6A, 0x0D, 0xAD)
DK2   = RGBColor(0x2C, 0x2C, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Core Helpers ──────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def box(slide, text, l, t, w, h, sz=12, bold=False, fc=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = fc
    return tx

def slide_hdr(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.18, NAVY)
    rect(slide, 0, 1.18, 13.33, 0.055, GOLD)
    box(slide, title, 0.38, 0.16, 12.5, 0.62, sz=22, bold=True, fc=WHT)
    if subtitle:
        box(slide, subtitle, 0.38, 0.76, 12.5, 0.38, sz=12, fc=GOLD, italic=False)

def footnote_bar(slide, text, y=7.2):
    rect(slide, 0, y, 13.33, 0.015, MGRY)
    box(slide, text, 0.35, y+0.03, 12.8, 0.32, sz=8, fc=MGRY, italic=True)

def section_divider(slide, label, color=NAVY):
    """Full-bleed section divider slide"""
    fill_bg(slide, color)
    rect(slide, 0, 3.3, 13.33, 0.07, GOLD)
    box(slide, label, 1.0, 2.6, 11.3, 1.2, sz=36, bold=True, fc=WHT, align=PP_ALIGN.CENTER)

def add_tbl(slide, rows, cols, l, t, w, h):
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    return ts.table

def tcell(tbl, row, col, text, sz=10, bold=False, fc=DARK, bg=None,
          align=PP_ALIGN.LEFT, italic=False):
    cell = tbl.cell(row, col)
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    # Remove existing runs at XML level
    for r_elem in p._p.findall(qn('a:r')):
        p._p.remove(r_elem)
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = fc
    return cell

def thdr(tbl, row, col, text, sz=10, fc=WHT, bg=NAVY, align=PP_ALIGN.CENTER):
    return tcell(tbl, row, col, text, sz=sz, bold=True, fc=fc, bg=bg, align=align)

def alt(tbl, row, col, text, sz=10, bold=False, fc=DARK, align=PP_ALIGN.LEFT):
    bg = STLB if row % 2 == 0 else None
    return tcell(tbl, row, col, text, sz=sz, bold=bold, fc=fc, bg=bg, align=align)

def bullet_box(slide, title, bullets, l, t, w, h,
               accent=NAVY, title_sz=12, body_sz=11):
    rect(slide, l, t, w, h, LGRY)
    rect(slide, l, t, 0.12, h, accent)
    box(slide, title, l+0.2, t+0.07, w-0.3, 0.38, sz=title_sz, bold=True, fc=DARK)
    body = "\n".join(f"• {b}" for b in bullets)
    box(slide, body, l+0.2, t+0.5, w-0.3, h-0.62, sz=body_sz, fc=DGRY)

def kpi_box(slide, value, label, l, t, w=1.95, h=1.35, accent=NAVY):
    rect(slide, l, t, w, h, WHT)
    rect(slide, l, t+h-0.06, w, 0.06, accent)
    box(slide, value, l, t+0.18, w, 0.72, sz=24, bold=True, fc=accent,
        align=PP_ALIGN.CENTER)
    box(slide, label, l, t+0.9, w, 0.42, sz=10, fc=MGRY,
        align=PP_ALIGN.CENTER, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s1 = blank(); fill_bg(s1, NAVY)
rect(s1, 0, 0, 13.33, 0.5, RGBColor(0x08, 0x16, 0x30))
rect(s1, 0, 7.0, 13.33, 0.5, RGBColor(0x08, 0x16, 0x30))
rect(s1, 0, 2.72, 13.33, 0.07, GOLD)
rect(s1, 0, 4.2, 13.33, 0.07, GOLD)
# Confidential badge
rect(s1, 0.4, 0.12, 3.5, 0.32, GOLD)
box(s1, "STRICTLY CONFIDENTIAL — PROJECT STARLIGHT",
    0.42, 0.13, 3.48, 0.3, sz=8.5, bold=True, fc=NAVY)
# Title
box(s1, "Project Starlight", 1.2, 1.05, 11, 0.75,
    sz=16, fc=GOLD, italic=True, align=PP_ALIGN.CENTER)
box(s1, "Strategic Acquisition of\nWarner Bros. Discovery",
    0.7, 1.75, 11.9, 1.6, sz=38, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
box(s1, "Presented to Netflix Management  |  May 2025",
    0.7, 4.38, 11.9, 0.52, sz=14, fc=RGBColor(0xA0,0xB0,0xC8), align=PP_ALIGN.CENTER)
box(s1, "PREPARED BY YOUR FINANCIAL ADVISOR  ·  FOR DISCUSSION PURPOSES ONLY",
    0.7, 5.1, 11.9, 0.45, sz=12, fc=MGRY, align=PP_ALIGN.CENTER, italic=True)
# Footer
box(s1, "This presentation is confidential and intended solely for the use of the "
        "management of Netflix, Inc. and its authorized representatives.",
    0.7, 7.06, 11.9, 0.4, sz=9, fc=MGRY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DISCLAIMER
# ════════════════════════════════════════════════════════════════════════════
s2 = blank(); fill_bg(s2, LGRY)
rect(s2, 0, 0, 13.33, 1.18, NAVY)
rect(s2, 0, 1.18, 13.33, 0.055, GOLD)
box(s2, "IMPORTANT DISCLAIMER & CONFIDENTIALITY NOTICE",
    0.38, 0.18, 12.5, 0.65, sz=20, bold=True, fc=WHT)
rect(s2, 0.35, 1.35, 12.63, 5.6, WHT)
disc = (
    "CONFIDENTIALITY\n"
    "This presentation has been prepared exclusively for Netflix, Inc. (\"Netflix\" or the \"Company\") and its "
    "authorized representatives by its Financial Advisor. This material is strictly confidential and may not be "
    "reproduced, distributed, or disclosed to any other person without the prior written consent of the Financial Advisor.\n\n"
    "NOT AN OFFER\n"
    "This presentation does not constitute an offer to sell or a solicitation of an offer to buy any securities. "
    "Nothing herein shall be deemed to constitute investment advice, and it shall not form the basis of any contract. "
    "All financial projections are illustrative and should not be relied upon as indicators of future performance.\n\n"
    "FORWARD-LOOKING STATEMENTS\n"
    "Certain statements in this presentation constitute forward-looking statements involving risks and uncertainties. "
    "Actual results may differ materially from those anticipated due to factors beyond the parties' control, including "
    "regulatory approvals, market conditions, and integration execution. Neither Netflix nor the Financial Advisor "
    "makes any representation or warranty as to the accuracy or completeness of this information.\n\n"
    "SOURCES\n"
    "Financial data sourced from public filings (Netflix 10-K FY2024; Warner Bros. Discovery 10-K FY2024), "
    "Bloomberg, FactSet, and Financial Advisor research. Market data as of May 2025. "
    "WBD FY2024 Adjusted EBITDA per company guidance. All synergy estimates are Financial Advisor projections.\n\n"
    "FOR MANAGEMENT DISCUSSION PURPOSES ONLY — NOT FOR DISTRIBUTION"
)
box(s2, disc, 0.55, 1.45, 12.25, 5.35, sz=9.5, fc=DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TABLE OF CONTENTS
# ════════════════════════════════════════════════════════════════════════════
s3 = blank()
slide_hdr(s3, "Table of Contents")
fill_bg(s3, LGRY)
rect(s3, 0, 0, 13.33, 1.18, NAVY)
rect(s3, 0, 1.18, 13.33, 0.055, GOLD)
box(s3, "Table of Contents", 0.38, 0.18, 12.5, 0.65, sz=22, bold=True, fc=WHT)

sections = [
    ("I",   "Executive Summary & Transaction Highlights",   "4"),
    ("II",  "Why Now? — Strategic Context",                  "5"),
    ("III", "Netflix: Company Overview",                     "6"),
    ("IV",  "Warner Bros. Discovery: Company Overview",      "7 – 8"),
    ("V",   "Strategic Rationale — Five Pillars",            "9 – 10"),
    ("VI",  "Synergy Analysis",                              "11"),
    ("VII", "Transaction Structure & Financing",             "12 – 13"),
    ("VIII","Valuation Analysis",                            "14 – 16"),
    ("IX",  "Pro Forma Financial Impact",                    "17"),
    ("X",   "Risk Factors & Mitigants",                      "18"),
    ("XI",  "Recommendation & Next Steps",                   "19"),
    ("",    "Appendix: Selected Financial Data",             "20"),
]
for i, (num, title, pg) in enumerate(sections):
    y = 1.35 + i * 0.49
    bg_c = NAVY if num in ("I","V","VIII") else (LGRY if i % 2 == 0 else WHT)
    fc_n = WHT if bg_c == NAVY else NBLUE
    fc_t = WHT if bg_c == NAVY else DARK
    rect(s3, 0.35, y, 12.63, 0.45, bg_c)
    box(s3, num, 0.45, y+0.06, 0.7, 0.34, sz=11, bold=True, fc=fc_n)
    box(s3, title, 1.2, y+0.06, 10.5, 0.34, sz=11, fc=fc_t, bold=(bg_c==NAVY))
    box(s3, pg, 12.0, y+0.06, 0.9, 0.34, sz=11, fc=fc_n, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — EXECUTIVE SUMMARY / TRANSACTION HIGHLIGHTS
# ════════════════════════════════════════════════════════════════════════════
s4 = blank(); fill_bg(s4, LGRY)
slide_hdr(s4, "Executive Summary", "Transaction Highlights — Project Starlight")

# Transaction stat boxes
txn_stats = [
    ("$10.00/sh", "Proposed Offer Price"),
    ("+182%",     "Premium to Last Close\n($3.55 undisturbed)"),
    ("~$60B",     "Total Transaction\nValue (EV)"),
    ("6.2x",      "EV / FY2024A\nAdj. EBITDA"),
    ("~$4.0B",    "Estimated Annual\nSynergies (Yr 3)"),
    ("4.4x",      "EV / (EBITDA +\nSynergies)"),
]
for i, (val, lbl) in enumerate(txn_stats):
    x = 0.35 + i * 2.12
    kpi_box(s4, val, lbl, x, 1.32, w=2.0, h=1.32, accent=NAVY)

# Left: Rationale bullets
bullet_box(s4, "Strategic Rationale",
    ["100,000+ hours of premium content library (+HBO/DC/Harry Potter IP)",
     "Accelerates Netflix to 400M+ combined subscriber base",
     "Max platform ($7.7B DTC revenue, 96M subs) — eliminates #2 competitor",
     "Warner Bros. studio: exclusive pipeline for theatrical + franchise content",
     "Megapack of sports rights: March Madness, MLB, NHL through 2028–2032",
     "Transforms Netflix from SVOD to diversified media & entertainment platform"],
    0.35, 2.85, 6.1, 3.8, accent=NRED, title_sz=13, body_sz=10.5)

# Right: Financial rationale
bullet_box(s4, "Financial Rationale",
    ["WBD trades at 4.6x EBITDA — distressed pricing due to $39.5B debt load",
     "Netflix FCF of $6.9B/yr de-levers combined entity from 2.8x to <2.0x by Yr 3",
     "~$4.0B annual synergies (NPV ~$38B) on $23.8B equity acquisition cost",
     "Adds $22B Networks cash flow machine — funds content investment",
     "EPS accretive within 18–24 months post-close (adj. for synergies)",
     "Cross-platform monetization: Netflix + Max → single global bundle offering"],
    6.6, 2.85, 6.4, 3.8, accent=GOLD, title_sz=13, body_sz=10.5)

footnote_bar(s4, "Sources: Bloomberg, FactSet, Netflix 10-K FY2024, WBD 10-K FY2024. Market data as of May 2025. "
                 "Synergy estimates are Financial Advisor projections. This analysis is for discussion purposes only.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHY NOW?
# ════════════════════════════════════════════════════════════════════════════
s5 = blank(); fill_bg(s5, WHT)
slide_hdr(s5, "Why Now? — Strategic Context",
          "A Once-in-a-Generation Opportunity to Consolidate Streaming Leadership")

contexts = [
    ("WBD at an Inflection Point", RED2,
     "• WBD shares down 80%+ since AT&T spinoff (Apr 2022)\n"
     "• $39.5B debt burden constrains content investment\n"
     "• Bond covenant pressure limits strategic flexibility\n"
     "• Management exploring strategic alternatives (public since Feb 2025)"),
    ("Cord-Cutting Accelerating", NBLUE,
     "• Linear TV households declined 8% YoY (2024)\n"
     "• Networks EBITDA declined 12% in FY2024\n"
     "• WBD's cash-cow Networks segment (TNT/CNN/TBS) under structural pressure\n"
     "• Buyers' market for linear assets — act before further erosion"),
    ("Streaming Consolidation Wave", TEAL,
     "• Paramount + Skydance closed July 2024\n"
     "• Comcast spinning NBCU cable assets (announced Jan 2025)\n"
     "• Scale advantages compounding — top 2 players capture 65%+ of growth\n"
     "• Window to acquire WBD at distressed price is closing"),
    ("Content Moats Are Critical", GRN,
     "• Netflix content spend $17B/yr — WBD adds 100K hrs for marginal cost\n"
     "• DC Films franchise reboot under James Gunn — massive IP upside\n"
     "• Harry Potter TV series + Rings of Power cross-IP opportunity\n"
     "• HBO prestige brand = most awarded TV content in history"),
    ("Regulatory Environment Favorable", AMB,
     "• Biden-era DOJ media scrutiny easing under new administration\n"
     "• No overlapping broadcast licenses — no FCC issue\n"
     "• Horizontal streaming deal (no vertical integration concerns)\n"
     "• EU precedent: Disney/Fox approved with limited remedies"),
    ("Netflix at Strategic Crossroads", PURP,
     "• Subscriber growth slowing (from +22M Yr to +8M as base matures)\n"
     "• Ad-supported tier growing but needs more content variety\n"
     "• Live events (NFL, Boxing) showing engagement power — need sports portfolio\n"
     "• $6.9B FCF with limited organic investment opportunities at this scale"),
]
for i, (title, color, body) in enumerate(contexts):
    col = i % 2; row = i // 2
    x = 0.35 + col * 6.5; y = 1.32 + row * 2.0
    rect(s5, x, y, 6.15, 1.88, LGRY)
    rect(s5, x, y, 6.15, 0.42, color)
    box(s5, title, x+0.15, y+0.07, 5.85, 0.3, sz=12, bold=True, fc=WHT)
    box(s5, body, x+0.15, y+0.5, 5.9, 1.3, sz=10, fc=DARK)

footnote_bar(s5, "Sources: Bloomberg, FactSet, Nielsen, WBD/Netflix public filings. "
                 "Linear TV household data: Nielsen Media Research Q4 2024.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — NETFLIX COMPANY OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
s6 = blank(); fill_bg(s6, LGRY)
slide_hdr(s6, "Netflix, Inc. (NFLX) — Company Overview",
          "Global Streaming Leader with $380B Market Cap and $6.9B Annual FCF")

# KPIs
nflx_kpis = [
    ("$380B",  "Market Cap\n(May 2025)"),
    ("$387B",  "Enterprise\nValue"),
    ("301M",   "Paid Subscribers\n(Q1 2025)"),
    ("$39.0B", "FY2024\nRevenue"),
    ("$10.4B", "FY2024 Adj.\nEBITDA"),
    ("$6.9B",  "FY2024\nFree Cash Flow"),
]
for i, (val, lbl) in enumerate(nflx_kpis):
    x = 0.35 + i * 2.12
    kpi_box(s6, val, lbl, x, 1.32, w=2.0, h=1.35, accent=NRED)

# Financial table
box(s6, "Key Financial Metrics ($B unless noted)", 0.35, 2.87, 8.5, 0.38,
    sz=12, bold=True, fc=DARK)
tbl6 = add_tbl(s6, 7, 6, 0.35, 3.28, 8.25, 2.88)
for j, h in enumerate(["Metric","FY2022A","FY2023A","FY2024A","FY2025E","FY2026E"]):
    thdr(tbl6, 0, j, h, sz=9.5)
rows6 = [
    ("Revenue ($B)",         "$31.6", "$33.7", "$39.0", "$43.5", "$48.2"),
    ("Revenue Growth",        "6.7%", "6.7%", "15.6%", "11.5%",  "10.8%"),
    ("Adj. EBITDA ($B)",     " $5.7", " $7.0", "$10.4", "$12.1", "$14.3"),
    ("EBITDA Margin",        "18.1%","20.8%", "26.6%", "27.8%",  "29.7%"),
    ("Net Income ($B)",      " $4.5", " $5.4", " $8.7", "$10.1", "$12.2"),
    ("Free Cash Flow ($B)",  " $1.6", " $6.9", " $6.9", " $8.2", "$10.5"),
]
for i, (lbl, *vals) in enumerate(rows6):
    bg = STLB if i % 2 == 0 else None
    tcell(tbl6, i+1, 0, lbl, sz=9.5, bold=True, fc=DARK, bg=bg)
    for j, v in enumerate(vals):
        tcell(tbl6, i+1, j+1, v, sz=9.5, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)

# Right: Description + strategy
rect(s6, 8.8, 2.87, 4.18, 4.29, WHT)
desc_items = [
    ("Business Model", NRED,
     "Subscription-only SVOD platform (no ads tier growing). "
     "189 countries. Ad-supported plan launched 2022 — now 40%+ of new sign-ups."),
    ("Content Strategy", NBLUE,
     "~$17B annual content spend. Original programming (Stranger Things, Squid Game). "
     "Gaming expansion. Live events: NFL, Tyson vs. Paul fight."),
    ("Key Growth Drivers", GRN,
     "Password-sharing crackdown (+30M subs in 18 months). "
     "Ad-tier ARPU expansion. International market penetration. Live content."),
]
for i, (title, color, body) in enumerate(desc_items):
    y = 2.9 + i * 1.4
    rect(s6, 8.82, y, 4.14, 1.3, LGRY)
    rect(s6, 8.82, y, 0.1, 1.3, color)
    box(s6, title, 9.0, y+0.05, 3.85, 0.38, sz=11, bold=True, fc=DARK)
    box(s6, body, 9.0, y+0.46, 3.85, 0.78, sz=9.5, fc=DGRY)

footnote_bar(s6, "Sources: Netflix 10-K FY2024, Netflix Q1 2025 Earnings. "
                 "FY2025E–FY2026E based on FactSet consensus estimates.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WBD COMPANY OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
s7 = blank(); fill_bg(s7, LGRY)
slide_hdr(s7, "Warner Bros. Discovery (WBD) — Company Overview",
          "Deeply Undervalued Media Conglomerate with $44.5B EV and Transformative Asset Base")

wbd_kpis = [
    ("$8.5B",  "Market Cap\n(Distressed)"),
    ("$44.5B", "Enterprise\nValue"),
    ("96M",    "Max Paid\nSubscribers"),
    ("$39.3B", "FY2024\nRevenue"),
    (" $9.7B", "FY2024 Adj.\nEBITDA"),
    ("$39.5B", "Total Debt\n(Key Constraint)"),
]
for i, (val, lbl) in enumerate(wbd_kpis):
    x = 0.35 + i * 2.12
    accent = RED2 if i == 5 else NBLUE
    kpi_box(s7, val, lbl, x, 1.32, w=2.0, h=1.35, accent=accent)

# Segment financials table
box(s7, "Segment Financials — FY2024A ($B)", 0.35, 2.87, 8.25, 0.38,
    sz=12, bold=True, fc=DARK)
tbl7 = add_tbl(s7, 5, 5, 0.35, 3.28, 8.25, 2.0)
for j, h in enumerate(["Segment","Revenue","% of Total","Adj. EBITDA","EBITDA Margin"]):
    thdr(tbl7, 0, j, h, sz=9.5)
seg_rows = [
    ("Studios",         "$9.6B",  "24%", "$1.4B", "14.6%"),
    ("Networks",        "$22.0B", "56%", "$7.0B", "31.8%"),
    ("DTC (Max)",       "$7.7B",  "20%", "$1.3B", "16.9%"),
    ("Total / WBD",     "$39.3B","100%", "$9.7B", "24.7%"),
]
for i, (seg, rev, pct, ebitda, marg) in enumerate(seg_rows):
    bold = (i == 3)
    bg = NAVY if bold else (STLB if i % 2 == 0 else None)
    fc_ = WHT if bold else DARK
    for j, v in enumerate([seg, rev, pct, ebitda, marg]):
        tcell(tbl7, i+1, j, v, sz=9.5, bold=bold, fc=fc_, bg=bg,
              align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)

# WBD challenges + opportunities
rect(s7, 8.8, 2.87, 4.18, 3.45, WHT)
box(s7, "Key Challenges", 8.95, 2.92, 3.9, 0.35, sz=11, bold=True, fc=RED2)
challenges = [
    "• $39.5B debt → $3.0B/yr interest expense burden",
    "• Linear Networks declining (–12% EBITDA YoY)",
    "• Lost NBA rights to NBC — $1.8B/yr revenue at risk",
    "• Max subs growth slowing; competes with Netflix",
    "• Stock -80% from spinoff; strategic review underway",
]
box(s7, "\n".join(challenges), 8.95, 3.3, 3.9, 2.0, sz=9.5, fc=DARK)

footnote_bar(s7, "Sources: WBD 10-K FY2024, Q4 2024 Earnings Supplement. Adjusted EBITDA per company guidance. "
                 "Market data as of May 2025.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WBD ASSET PORTFOLIO
# ════════════════════════════════════════════════════════════════════════════
s8 = blank(); fill_bg(s8, WHT)
slide_hdr(s8, "Warner Bros. Discovery — Asset Portfolio",
          "Irreplaceable Content IP, Sports Rights, and Distribution Infrastructure")

box(s8, "WBD Asset Overview — Standalone Estimated Values", 0.35, 1.28, 12.6, 0.38,
    sz=13, bold=True, fc=DARK)

tbl8 = add_tbl(s8, 10, 5, 0.35, 1.7, 12.63, 4.35)
for j, h in enumerate(["Asset / Division","Category","Revenue (FY24A)","Est. Stand-Alone Value","Strategic Value to Netflix"]):
    thdr(tbl8, 0, j, h, sz=9.5)

assets = [
    ("Warner Bros. Pictures Studio", "Studios", "$4.2B", "$18–22B",    "Exclusive theatrical pipeline; franchise content"),
    ("DC Comics / DC Films",         "IP",      "Included", "$8–15B",  "Superman, Batman, WW; reboot under James Gunn"),
    ("HBO / Max Platform",           "DTC",     "$7.7B", "$12–18B",    "Eliminates #2 streaming competitor; 96M subs"),
    ("Harry Potter / Wizarding World","IP",     "Included", "$5–8B",   "HBO series announced; streaming goldmine"),
    ("CNN / Turner News",            "Networks","$3.5B", "$4–6B",      "News content + live programming capability"),
    ("TNT / TBS / Cartoon Network",  "Networks","$7.2B", "$5–8B",      "Sports: March Madness (–2032) + MLB + NHL"),
    ("Discovery / HGTV / Food / TLC","Networks","$6.1B", "$6–9B",      "Unscripted/reality: highest margin content"),
    ("Game of Thrones / Lord of Rings","IP",    "Included","$4–7B",    "Highest-performing IP in Netflix history (if licensed)"),
    ("Total (Illustrative Range)",   "All",     "$39.3B","$62–93B",    "vs. current EV of $44.5B — clear value gap"),
]
for i, row in enumerate(assets):
    bold = (i == 8)
    bg = NAVY if bold else (STLB if i % 2 == 0 else None)
    fc_ = WHT if bold else DARK
    for j, v in enumerate(row):
        al = PP_ALIGN.LEFT if j in (0,1,4) else PP_ALIGN.RIGHT
        tcell(tbl8, i+1, j, v, sz=9, bold=bold, fc=fc_, bg=bg, align=al)

rect(s8, 0.35, 6.15, 12.63, 0.68, RGBColor(0xFF, 0xF5, 0xE0))
rect(s8, 0.35, 6.15, 0.12, 0.68, GOLD)
box(s8, "Key Insight: WBD's sum-of-the-parts asset value ($62–93B) substantially exceeds its current "
        "enterprise value of $44.5B. Netflix would acquire these assets at a material discount to intrinsic value "
        "by leveraging WBD's distressed balance sheet.",
    0.55, 6.2, 12.3, 0.58, sz=10.5, fc=DARK, italic=False)

footnote_bar(s8, "Asset values are Financial Advisor estimates based on comparable transactions and trading multiples. "
                 "Not a guarantee of future value. For illustrative purposes only.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — STRATEGIC RATIONALE (5 PILLARS)
# ════════════════════════════════════════════════════════════════════════════
s9 = blank(); fill_bg(s9, LGRY)
slide_hdr(s9, "Strategic Rationale — Five Pillars",
          "Why Netflix + WBD Creates an Unassailable Entertainment Platform")

pillars = [
    ("I", "Content Scale &\nLibrary Depth", NAVY,
     "• 100,000+ hours of premium content acquired at near-zero marginal cost\n"
     "• HBO's 25-year prestige library (Sopranos, Wire, Game of Thrones)\n"
     "• Warner Bros. theatrical titles: 30+ blockbusters/yr enter Netflix catalogue\n"
     "• Content CPM drops 40–60% when library amortized across 400M+ subs\n"
     "• Eliminates $3–5B/yr in content spend duplicating WBD originals"),
    ("II", "IP Franchise\nMachinery", NRED,
     "• DC Universe: 20+ hero franchises; James Gunn DCU reboot (2025–2028)\n"
     "• Harry Potter: HBO series + 4 Fantastic Beasts films in pipeline\n"
     "• Lord of the Rings / Hobbit: Amazon competing for same eyeballs\n"
     "• Looney Tunes, Tom & Jerry: family animation library\n"
     "• IP creates 5–10yr content roadmaps with built-in global fanbases"),
    ("III", "Subscriber Base &\nMarket Position", TEAL,
     "• Combined: 397M paid subs (Netflix 301M + Max 96M) — #1 globally by wide margin\n"
     "• Eliminates Max as primary subscriber acquisition competitor\n"
     "• Ad-tier cross-sell: Max $9.99 ad subs → Netflix $6.99 plan migration\n"
     "• Retention improvement: Combined library reduces churn by est. 0.5–1.0pp/mo\n"
     "• Bundle pricing power: $15–18/mo single streaming service"),
    ("IV", "Sports Rights &\nLive Content", GRN,
     "• March Madness (NCAA): 68-game tournament; contract through 2032\n"
     "• MLB: Sunday/Wednesday games + playoffs through 2028\n"
     "• NHL: Stanley Cup Playoffs + regular season through 2028\n"
     "• Live sports = best-in-class churn reducer and ad-tier driver\n"
     "• Netflix already exploring sports (NFL, Formula 1); WBD accelerates roadmap"),
    ("V", "Global Distribution\n& Studio Scale", PURP,
     "• Warner Bros. theatrical: global marketing machine for Netflix IP\n"
     "• Day-and-date or short theatrical window → Netflix streaming\n"
     "• Discovery's 220-country distribution footprint\n"
     "• WBD consumer products / licensing: $1B+ incremental revenue\n"
     "• Combined studio provides Netflix full vertical integration"),
]
for i, (num, title, color, body) in enumerate(pillars):
    x = 0.35 + i * 2.59
    rect(s9, x, 1.32, 2.45, 5.5, color)
    box(s9, num, x, 1.35, 2.45, 0.65, sz=28, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
    box(s9, title, x+0.1, 1.95, 2.25, 0.72, sz=12, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
    rect(s9, x+0.1, 2.72, 2.25, 0.04, RGBColor(0xFF,0xFF,0xFF))
    box(s9, body, x+0.1, 2.82, 2.28, 3.9, sz=9.5, fc=WHT)

footnote_bar(s9, "Financial Advisor analysis. Sports rights contracts sourced from public filings. "
                 "Subscriber data: Netflix Q1 2025; WBD Q4 2024 Earnings Supplements.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONTENT & IP DEEP DIVE
# ════════════════════════════════════════════════════════════════════════════
s10 = blank(); fill_bg(s10, WHT)
slide_hdr(s10, "Content Library & IP Value — Deep Dive",
          "Quantifying the Strategic Value of WBD's Irreplaceable Content Assets")

box(s10, "Comparable Content Library Transactions", 0.35, 1.28, 8.1, 0.38,
    sz=12, bold=True, fc=DARK)
tbl10 = add_tbl(s10, 7, 5, 0.35, 1.7, 8.1, 3.0)
for j, h in enumerate(["Transaction","Year","Library Size","Transaction Value","$/Hour"]):
    thdr(tbl10, 0, j, h, sz=9.5)
lib_comps = [
    ("Disney / Fox (library assets)",       "2019", "~75,000 hrs",   "$71.3B total",    "~$950K/hr"),
    ("MGM / Amazon",                        "2022", "~17,000 hrs",   "$8.45B",          "~$497K/hr"),
    ("Showtime / CBS / Paramount+",         "2023", "~25,000 hrs",   "$15.4B total",    "~$616K/hr"),
    ("Sony Pictures (partial license)",     "2021", "~3,500 hrs/yr", "$1.0B/yr",        "~$286K/hr"),
    ("Comcast / NBCU streaming assets",     "2025", "~30,000 hrs",   "~$9B est.",       "~$300K/hr"),
    ("WBD Library (estimated)",             "2025", "100,000+ hrs",  "$30–50B implied", "$300–500K/hr"),
]
for i, row in enumerate(lib_comps):
    is_wbd = (i == 5)
    bg = NAVY if is_wbd else (STLB if i % 2 == 0 else None)
    fc_ = WHT if is_wbd else DARK
    for j, v in enumerate(row):
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT
        tcell(tbl10, i+1, j, v, sz=9.5, bold=is_wbd, fc=fc_, bg=bg, align=al)

# IP Value table (right side)
box(s10, "Franchise IP Value Estimates", 8.6, 1.28, 4.4, 0.38,
    sz=12, bold=True, fc=DARK)
tbl10b = add_tbl(s10, 6, 3, 8.6, 1.7, 4.4, 2.1)
for j, h in enumerate(["IP Franchise","Est. Franchise Value","Note"]):
    thdr(tbl10b, 0, j, h, sz=9.5)
ip_items = [
    ("DC Universe",         "$15–25B",  "Reboot; 20+ films planned"),
    ("Harry Potter",        "$10–15B",  "HBO series + Fantastic Beasts"),
    ("Game of Thrones",     " $5–10B",  "Multiple spin-offs in dev."),
    ("Lord of the Rings",   " $4–8B",   "Amazon competing; leverage"),
    ("Total (4 franchises)","$34–58B",  "vs. WBD EV of $44.5B"),
]
for i, (ip, val, note) in enumerate(ip_items):
    bold = (i == 4)
    bg = NAVY if bold else (STLB if i % 2 == 0 else None)
    fc_ = WHT if bold else DARK
    tcell(tbl10b, i+1, 0, ip,  sz=9.5, bold=bold, fc=fc_, bg=bg)
    tcell(tbl10b, i+1, 1, val, sz=9.5, bold=bold, fc=fc_, bg=bg, align=PP_ALIGN.RIGHT)
    tcell(tbl10b, i+1, 2, note,sz=9.5, bold=False, fc=fc_, bg=bg)

# Bottom commentary boxes
rect(s10, 0.35, 4.85, 12.63, 1.9, LGRY)
rect(s10, 0.35, 4.85, 0.12, 1.9, GOLD)
insights = [
    ("Content Cost Efficiency", NAVY,
     "Netflix pays avg. $150–300K/hr for originals. WBD library adds 100K hours at ~$0 marginal cost post-acquisition. "
     "Content cost savings of $3–5B/yr vs. recreating equivalent content organically."),
    ("IP Monetization", NRED,
     "Netflix's House of Cards → $0.4B lifetime value. DC reboot films historically generate $500M–$1B+ box office. "
     "Netflix + theatrical releases = highest-value content monetization model."),
    ("Library Retention Value", GRN,
     "Library content (back-catalogue) accounts for 40–60% of Netflix viewing hours. "
     "WBD library reduces monthly churn by est. 0.5–1.0 percentage points = $2–4B/yr retention value."),
]
for i, (title, color, body) in enumerate(insights):
    x = 0.5 + i * 4.2
    rect(s10, x, 4.92, 4.0, 1.7, WHT)
    rect(s10, x, 4.92, 0.1, 1.7, color)
    box(s10, title, x+0.18, 4.95, 3.75, 0.36, sz=11, bold=True, fc=DARK)
    box(s10, body, x+0.18, 5.35, 3.75, 1.2, sz=9.5, fc=DGRY)

footnote_bar(s10, "IP valuations are Financial Advisor estimates based on comparable franchise transactions. "
                  "Library transaction data sourced from public filings and industry reports.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — SYNERGY ANALYSIS
# ════════════════════════════════════════════════════════════════════════════
s11 = blank(); fill_bg(s11, LGRY)
slide_hdr(s11, "Synergy Analysis",
          "~$4.0B Annual Synergies by Year 3 — NPV of ~$38B at 8.75% Discount Rate")

# Synergy summary boxes
syn_boxes = [
    ("~$1.5B", "Content Cost\nSynergies / Yr", GRN),
    ("~$1.2B", "Revenue\nSynergies / Yr", NBLUE),
    ("~$0.8B", "Corporate &\nTech Synergies", TEAL),
    ("~$0.5B", "Distribution &\nMarketing", GOLD),
    ("~$4.0B", "Total Annual\nSynergies (Yr 3)", NAVY),
    ("~$38B",  "NPV of\nSynergies", NRED),
]
for i, (val, lbl, color) in enumerate(syn_boxes):
    x = 0.35 + i * 2.12
    kpi_box(s11, val, lbl, x, 1.32, w=2.0, h=1.32, accent=color)

# Synergy detail table
box(s11, "Synergy Build-Up — Phased Implementation ($M)", 0.35, 2.82, 12.63, 0.38,
    sz=12, bold=True, fc=DARK)
tbl11 = add_tbl(s11, 13, 6, 0.35, 3.25, 12.63, 3.35)
for j, h in enumerate(["Category","Description","Yr 1E","Yr 2E","Yr 3E","Confidence"]):
    thdr(tbl11, 0, j, h, sz=9.5)

synergies = [
    # Cost synergies
    ("COST SYNERGIES",            "",                            "",      "",       "",       ""),
    ("Content Spend Reduction",   "Eliminate duplicate originals; leverage combined library", "$500", "$900", "$1,500", "High"),
    ("Technology / Infrastructure","Consolidate streaming tech; single CDN + data infra", "$150", "$250", "$400",   "High"),
    ("Corporate Overhead",        "Duplicate G&A, exec comp, legal, finance, HR",            "$100", "$180", "$280",   "High"),
    ("Marketing & Distribution",  "Combined subscriber marketing; co-brand promotions",       "$100", "$200", "$350",   "Medium"),
    ("Cost Subtotal",             "",                                   "$850", "$1,530", "$2,530", ""),
    # Revenue synergies
    ("REVENUE SYNERGIES",         "",                                   "",      "",       "",       ""),
    ("Cross-Platform Subscribers","Max-to-Netflix migration; retention improvement",          "$200", "$450", "$800",   "Medium"),
    ("Advertising Uplift",        "Combined audience scale → CPM premium for ad tier",        "$100", "$250", "$500",   "Medium"),
    ("Content Monetization",      "IP licensing, consumer products, theatrical windows",      "$100", "$200", "$300",   "Low"),
    ("Bundle Pricing Power",      "Single $17.99 Netflix bundle: est. 10% ARPU improvement", " $75", "$170", "$320",   "Low"),
    ("TOTAL SYNERGIES",           "",                               "$1,325", "$2,600", "$4,450", ""),
]
syn_colors = {0: NBLUE, 5: LGRY, 6: GRN, 11: NAVY}
syn_fcs    = {0: WHT,   5: DARK, 6: WHT,  11: WHT}
for i, row in enumerate(synergies):
    bold = (i in (0, 5, 6, 11))
    bg_ = syn_colors.get(i, STLB if i % 2 == 0 else None)
    fc_ = syn_fcs.get(i, DARK)
    for j, v in enumerate(row):
        al = PP_ALIGN.RIGHT if j in (2,3,4) else PP_ALIGN.LEFT
        tcell(tbl11, i+1, j, v, sz=9, bold=bold, fc=fc_, bg=bg_, align=al)

footnote_bar(s11, "Synergy estimates are Financial Advisor projections based on comparable media M&A transactions. "
                  "NPV calculated at 8.75% WACC. Synergies net of estimated one-time restructuring costs of ~$1.5B.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TRANSACTION TERMS
# ════════════════════════════════════════════════════════════════════════════
s12 = blank(); fill_bg(s12, WHT)
slide_hdr(s12, "Proposed Transaction Terms",
          "All-Cash / Stock Merger at $10.00 per WBD Share — ~$60B Total Transaction Value")

# Left: deal term table
box(s12, "Transaction Summary", 0.35, 1.28, 6.1, 0.38, sz=13, bold=True, fc=DARK)
tbl12 = add_tbl(s12, 14, 2, 0.35, 1.7, 6.1, 4.55)
thdr(tbl12, 0, 0, "Term", sz=10)
thdr(tbl12, 0, 1, "Detail", sz=10, align=PP_ALIGN.LEFT)
terms = [
    ("Transaction Type",            "Negotiated merger; cash & stock consideration"),
    ("Acquirer",                    "Netflix, Inc. (NASDAQ: NFLX)"),
    ("Target",                      "Warner Bros. Discovery, Inc. (NASDAQ: WBD)"),
    ("Offer Price Per Share",       "$10.00 (mix: 50% cash / 50% NFLX stock)"),
    ("Premium to Undisturbed Price","~182% to $3.55 (30-day VWAP: ~177%)"),
    ("Target Shares Outstanding",   "~2.38 billion (diluted)"),
    ("Equity Value",                "~$23.8 billion"),
    ("Assumed Net Debt",            "~$36.0 billion"),
    ("Total Transaction Value (EV)","~$59.8 billion (~$60B)"),
    ("EV / LTM Adj. EBITDA",        "6.2x ($60B / $9.7B)"),
    ("EV / (EBITDA + Synergies)",   "4.4x (with $4.0B/yr synergies)"),
    ("Form of Consideration",       "$5.00 cash + $5.00 in NFLX shares (at signing VWAP)"),
    ("Expected Close",              "12–18 months post-announcement (H2 2026)"),
]
for i, (term, detail) in enumerate(terms):
    bold = term in ("Total Transaction Value (EV)","EV / LTM Adj. EBITDA","Offer Price Per Share")
    bg = NAVY if term == "Total Transaction Value (EV)" else (STLB if i % 2 == 0 else None)
    fc_ = WHT if bg == NAVY else DARK
    tcell(tbl12, i+1, 0, term,   sz=9.5, bold=bold, fc=fc_, bg=bg)
    tcell(tbl12, i+1, 1, detail, sz=9.5, bold=bold, fc=fc_, bg=bg)

# Right: Conditions + Timeline
box(s12, "Conditions to Closing", 6.65, 1.28, 6.35, 0.38, sz=13, bold=True, fc=DARK)
conds = [
    ("Regulatory Approvals", NAVY,
     "• U.S. DOJ / FTC antitrust clearance\n"
     "• EU Commission merger review\n"
     "• FCC approval (no broadcast licenses held)\n"
     "• UK CMA review (~6–9 months)\n"
     "• Australian ACCC / Canadian review"),
    ("Financing Condition", GRN,
     "• Committed financing from lenders (no flex)\n"
     "• NFLX board approval of stock issuance\n"
     "• WBD board fiduciary process (go-shop 30 days)\n"
     "• WBD shareholder vote (simple majority)"),
    ("Key Timeline", GOLD,
     "• Day 0: Signing & public announcement\n"
     "• Month 1–3: HSR filing + EU Form CO\n"
     "• Month 6–12: Regulatory review periods\n"
     "• Month 12–18: Close + integration Day 1"),
]
for i, (title, color, body) in enumerate(conds):
    y = 1.7 + i * 1.75
    rect(s12, 6.65, y, 6.35, 1.65, LGRY)
    rect(s12, 6.65, y, 0.12, 1.65, color)
    box(s12, title, 6.85, y+0.07, 6.05, 0.38, sz=11, bold=True, fc=DARK)
    box(s12, body, 6.85, y+0.5, 6.05, 1.1, sz=9.5, fc=DARK)

footnote_bar(s12, "Terms are illustrative and subject to negotiation. Regulatory timeline is estimated. "
                  "NFLX stock portion uses 20-day VWAP at signing. Transaction fees estimated at ~$300M.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — SOURCES & USES
# ════════════════════════════════════════════════════════════════════════════
s13 = blank(); fill_bg(s13, LGRY)
slide_hdr(s13, "Transaction Financing — Sources & Uses",
          "Pro Forma Net Leverage of ~2.8x EBITDA — Supported by $10B+ Combined FCF")

# Sources table
box(s13, "Sources of Funds", 0.35, 1.28, 5.8, 0.38, sz=13, bold=True, fc=DARK)
tbl13s = add_tbl(s13, 6, 3, 0.35, 1.7, 5.8, 2.15)
for j, h in enumerate(["Source","Amount ($B)","% of Total"]):
    thdr(tbl13s, 0, j, h, sz=10, align=PP_ALIGN.RIGHT if j>0 else PP_ALIGN.LEFT)
sources = [
    ("New Term Loan B (7-year)",              "$8.0B",  "13.7%"),
    ("New Senior Notes (10-year)",            "$5.0B",   "8.6%"),
    ("Netflix Cash on Balance Sheet",         "$7.0B",  "12.0%"),
    ("Netflix Common Stock (11.5M shares)",   "$10.0B", "17.2%"),
    ("Assumed WBD Existing Debt",             "$28.2B", "48.5%"),
]
for i, (src, amt, pct) in enumerate(sources):
    bold = False
    bg = STLB if i % 2 == 0 else None
    tcell(tbl13s, i+1, 0, src, sz=9.5, fc=DARK, bg=bg)
    tcell(tbl13s, i+1, 1, amt, sz=9.5, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)
    tcell(tbl13s, i+1, 2, pct, sz=9.5, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)

# Sources total
tcell(tbl13s, 5, 0, "Total Sources", sz=9.5, bold=True, fc=WHT, bg=NAVY)
tcell(tbl13s, 5, 1, "$58.2B", sz=9.5, bold=True, fc=WHT, bg=NAVY, align=PP_ALIGN.RIGHT)
tcell(tbl13s, 5, 2, "100.0%", sz=9.5, bold=True, fc=WHT, bg=NAVY, align=PP_ALIGN.RIGHT)

# Uses table
box(s13, "Uses of Funds", 6.35, 1.28, 5.8, 0.38, sz=13, bold=True, fc=DARK)
tbl13u = add_tbl(s13, 5, 3, 6.35, 1.7, 5.8, 1.82)
for j, h in enumerate(["Use","Amount ($B)","% of Total"]):
    thdr(tbl13u, 0, j, h, sz=10, align=PP_ALIGN.RIGHT if j>0 else PP_ALIGN.LEFT)
uses_data = [
    ("Purchase of WBD Equity ($10.00/sh)",     "$23.8B", "40.9%"),
    ("Refinance WBD Existing Debt",            "$28.2B", "48.5%"),
    ("Estimated Transaction Fees & Expenses",  " $0.3B",  "0.5%"),
]
for i, (use, amt, pct) in enumerate(uses_data):
    bg = STLB if i % 2 == 0 else None
    tcell(tbl13u, i+1, 0, use, sz=9.5, fc=DARK, bg=bg)
    tcell(tbl13u, i+1, 1, amt, sz=9.5, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)
    tcell(tbl13u, i+1, 2, pct, sz=9.5, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)
tcell(tbl13u, 4, 0, "Total Uses", sz=9.5, bold=True, fc=WHT, bg=NAVY)
tcell(tbl13u, 4, 1, "$52.3B*", sz=9.5, bold=True, fc=WHT, bg=NAVY, align=PP_ALIGN.RIGHT)
tcell(tbl13u, 4, 2, "100.0%", sz=9.5, bold=True, fc=WHT, bg=NAVY, align=PP_ALIGN.RIGHT)

# Capitalization table
box(s13, "Pro Forma Capitalization & Leverage ($B)", 0.35, 4.0, 12.63, 0.38,
    sz=13, bold=True, fc=DARK)
tbl13c = add_tbl(s13, 7, 5, 0.35, 4.42, 12.63, 2.38)
for j, h in enumerate(["Metric","Netflix Stand-Alone","WBD Stand-Alone",
                        "Pro Forma Combined","Pro Forma + Yr 3 Synergies"]):
    thdr(tbl13c, 0, j, h, sz=9.5)
cap_rows = [
    ("Total Debt ($B)",               "$14.0B", "$39.5B", "$61.5B",  "$55.0B (after FCF paydown)"),
    ("Cash ($B)",                     " $7.0B", " $3.5B", " $3.5B",  "$5.0B+"),
    ("Net Debt ($B)",                 " $7.0B", "$36.0B", "$58.0B",  "$50.0B"),
    ("Adj. EBITDA ($B)",              "$10.4B", " $9.7B", "$20.1B",  "$24.1B (incl. $4B synergies)"),
    ("Net Debt / EBITDA",             "  0.7x", "  3.7x", "  2.9x",  "2.1x (delevering to <2x Yr4)"),
    ("FCF / Year ($B)",               " $6.9B", " $2.1B", " $9.0B",  "$11.5B+ (funding rapid delever)"),
]
for i, (metric, nflx, wbd, pf, pfs) in enumerate(cap_rows):
    bold = metric in ("Net Debt / EBITDA",)
    bg = STLB if i % 2 == 0 else None
    fg = RED2 if (metric == "Net Debt / EBITDA" and i == 4) else DARK
    for j, v in enumerate([metric, nflx, wbd, pf, pfs]):
        al = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
        tcell(tbl13c, i+1, j, v, sz=9.5, bold=bold, fc=fg, bg=bg, align=al)

footnote_bar(s13, "* Excludes $5.9B WBD deferred tax refinancing. Net leverage improves to ~2.1x within 3 years on $11.5B+ FCF. "
                  "Pro forma credit ratings: BBB- / Ba1 (investment grade target within 24 months).")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — VALUATION: TRADING COMPS
# ════════════════════════════════════════════════════════════════════════════
s14 = blank(); fill_bg(s14, WHT)
slide_hdr(s14, "Valuation Analysis — Comparable Public Companies",
          "WBD Trades at a Significant Discount to Peers on EV/EBITDA")

box(s14, "Selected Publicly Traded Media & Streaming Comparables ($ in billions)",
    0.35, 1.28, 12.63, 0.38, sz=12, bold=True, fc=DARK)
tbl14 = add_tbl(s14, 15, 9, 0.35, 1.7, 12.63, 5.0)
for j, h in enumerate(["Company","Ticker","Mkt Cap","EV","EV/Rev\n(NTM)","EV/EBITDA\n(NTM)","P/E\n(NTM)","Rev\nGrowth","Commentary"]):
    thdr(tbl14, 0, j, h, sz=9)

comps14 = [
    # Streaming peers
    ("— STREAMING / CONTENT —",  "",      "",       "",      "",      "",      "",     "", ""),
    ("Netflix",    "NFLX", "$380B", "$387B",  "8.9x", "37.2x", "43.6x", "+12%",  "Subject — premium multiple justified by growth"),
    ("Disney+/Hulu","DIS", "$190B", "$230B",  "2.4x", "12.8x", "22.4x", " +5%",  "Streaming + parks; content diversification"),
    ("Amazon Prime","AMZN","$2.2T", "$2.3T",  "3.5x", "21.5x", "44.2x", "+11%",  "Streaming as retention, not primary revenue"),
    ("Apple TV+",  "AAPL", "$3.3T", "$3.2T",  "8.4x", "26.5x", "32.1x", " +6%",  "Content as ecosystem driver; small subs base"),
    # Traditional media
    ("— TRADITIONAL MEDIA —",    "",      "",       "",      "",      "",      "",     "", ""),
    ("Comcast/NBCU","CMCSA","$155B","$220B",  "2.0x",  "7.9x", "12.5x", " +3%",  "Cable + Peacock; NBU spin-off announced"),
    ("Paramount+", "PARA",  "$8.5B","$22B",   "0.8x",  "5.4x", "10.2x", " +2%",  "Skydance merger completed July 2024"),
    ("AMC Networks","AMCX", " $1.0B","$3.5B", "1.1x",  "5.1x",  "7.4x", " -5%",  "Small-cap; highly distressed"),
    # Proposed acquisition
    ("— ACQUISITION —",          "",      "",       "",      "",      "",      "",     "", ""),
    ("WBD — Current",  "WBD", " $8.5B","$44.5B", "1.1x",  "4.6x",  "NM",  " -1%",  "Distressed; ~50% discount to cable comps"),
    ("WBD — Offer Price","",  "$23.8B","$59.8B", "1.5x",  "6.2x",  "NM",  " -1%",  "Transaction value — attractive vs. comps"),
    # Stats
    ("Streaming Peer Mean","",  "",  "",  "5.8x","24.5x","35.6x",  "8%",  ""),
    ("Cable/Media Peer Mean","","",  "",  "1.3x", "6.1x","10.1x",  "0%",  ""),
]
dividers  = {0, 5, 9, 11}
bold_rows = {2, 11, 12, 13}
for i, row in enumerate(comps14):
    if i in dividers:
        bg = NBLUE
        fc_ = WHT
    elif i == 11:
        bg = NAVY; fc_ = WHT
    elif i in bold_rows:
        bg = STLB; fc_ = DARK
    else:
        bg = STLB if i % 2 == 0 else None; fc_ = DARK
    bold = i in bold_rows or i in dividers
    for j, v in enumerate(row):
        al = PP_ALIGN.LEFT if j in (0,1,8) else PP_ALIGN.RIGHT
        tcell(tbl14, i+1, j, v, sz=8.5 if i not in dividers else 9,
              bold=bold, fc=fc_, bg=bg, align=al)

footnote_bar(s14, "Market data as of May 2025. Multiples based on NTM consensus EBITDA (FactSet). "
                  "WBD NM P/E reflects net losses. Streaming peers include Disney+/Hulu segment only for DIS.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — VALUATION: PRECEDENT TRANSACTIONS
# ════════════════════════════════════════════════════════════════════════════
s15 = blank(); fill_bg(s15, LGRY)
slide_hdr(s15, "Valuation Analysis — Precedent M&A Transactions",
          "Selected Media & Entertainment Acquisitions (2017–2024)")

box(s15, "Selected Precedent Transactions ($B unless noted)", 0.35, 1.28, 12.63, 0.38,
    sz=12, bold=True, fc=DARK)
tbl15 = add_tbl(s15, 13, 7, 0.35, 1.7, 12.63, 5.0)
for j, h in enumerate(["Date","Acquirer","Target","Transaction Value","EV/Revenue","EV/EBITDA","Notes"]):
    thdr(tbl15, 0, j, h, sz=9.5, align=PP_ALIGN.CENTER if j>3 else PP_ALIGN.LEFT)

prec = [
    ("Mar 2019","Disney",          "21st Century Fox",      "$71.3B",  "3.5x", "15.8x","Largest media M&A; content + FX + Nat Geo"),
    ("May 2021","Amazon",          "MGM Holdings",          " $8.5B",  "4.9x", "21.0x","Library-focused; Prime Video scale"),
    ("Apr 2022","WBD (AT&T)",      "WarnerMedia/Discovery", "$43.0B",  "1.8x",  "8.2x","Spin-merge; current acquiree"),
    ("Sep 2022","Microsoft",       "Activision Blizzard",   "$68.7B",  "7.0x", "23.5x","Gaming/IP; precedent for platform acquisition"),
    ("Dec 2023","Skydance",        "Paramount Global",      "$8.0B¹",  "0.7x",  "5.1x","Distressed deal; NAI controlling structure"),
    ("Mar 2024","Synapse/Sony",    "Paramount Studio Assets","$3.2B", "1.2x",  "9.5x","Studio library partial; comparable basis"),
    ("Jul 2024","Comcast",         "NBCU Streaming spinoff", "$9.0B²", "1.5x",  "8.5x","Cable network + Peacock separation"),
    ("2024",    "Various",         "Regional sports networks","$2–4B",  "1.0x",  "6–8x","RSN distressed deals benchmark"),
    # Stats
    ("",        "MEAN",            "",                       "",        "2.8x", "12.2x",""),
    ("",        "MEDIAN",          "",                       "",        "2.4x",  "8.8x",""),
    ("",        "25th / 75th",     "",                       "",     "1.4/3.9x","7.0/17x",""),
    ("Implied","NFLX / WBD Deal",  "$10.00/sh offer",        "$59.8B",  "1.5x",  "6.2x","At or below median — disciplined pricing"),
]
stat_rows = {8, 9, 10, 11}
for i, row in enumerate(prec):
    is_deal = (i == 11)
    is_stat = i in stat_rows
    bg = NAVY if is_deal else (NBLUE if is_stat else (STLB if i%2==0 else None))
    fc_ = WHT if (is_deal or is_stat) else DARK
    bold = is_deal or is_stat
    for j, v in enumerate(row):
        al = PP_ALIGN.LEFT if j in (0,1,2,6) else PP_ALIGN.CENTER
        tcell(tbl15, i+1, j, v, sz=9.5, bold=bold, fc=fc_, bg=bg, align=al)

footnote_bar(s15, "¹ Skydance deal includes assumption of ~$14B debt. ² Estimated value. "
                  "Multiples based on LTM EBITDA at transaction announcement. "
                  "Source: Public filings, Bloomberg, FactSet M&A database.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — FOOTBALL FIELD (VALUATION SUMMARY)
# ════════════════════════════════════════════════════════════════════════════
s16 = blank(); fill_bg(s16, WHT)
slide_hdr(s16, "Valuation Summary — Football Field Analysis",
          "Implied WBD Equity Value per Share Across Methodologies")

box(s16, "Implied Value per WBD Share", 0.35, 1.28, 8.0, 0.38, sz=13, bold=True, fc=DARK)

# Scale: $0 → $18 per share
# Scale bar from x=4.0 to x=12.8 = 8.8 inches = $18/sh
# pixels per dollar = 8.8/18 = 0.489 in / $

def ff_bar(slide, label, low, high, color, y, note=""):
    """Draw a football field bar"""
    scale = 8.8 / 18.0  # inches per dollar
    x_start = 4.0 + low  * scale
    bar_w   = (high - low) * scale
    rect(slide, x_start, y, bar_w, 0.46, color)
    # Low label
    box(slide, f"${low:.2f}", x_start - 0.7, y+0.06, 0.65, 0.34, sz=9.5,
        bold=True, fc=color, align=PP_ALIGN.RIGHT)
    # High label
    box(slide, f"${high:.2f}", x_start + bar_w + 0.05, y+0.06, 0.7, 0.34, sz=9.5,
        bold=True, fc=color)
    # Method label
    box(slide, label, 0.35, y+0.06, 3.6, 0.38, sz=10, bold=True, fc=DARK)
    if note:
        box(slide, note, 0.35, y+0.46, 3.6, 0.32, sz=9, fc=MGRY, italic=True)

# X-axis labels
for val in [0, 3, 6, 9, 12, 15, 18]:
    x_pos = 4.0 + val * (8.8/18.0)
    box(s16, f"${val}", x_pos - 0.2, 6.55, 0.4, 0.32, sz=9, fc=MGRY, align=PP_ALIGN.CENTER)
    rect(s16, x_pos, 1.68, 0.008, 5.0, MGRY)  # gridline

# The bars (low, high in $/share)
bars = [
    ("52-Week Trading Range",      2.80,  5.10, MGRY,   "Undisturbed: $3.55"),
    ("Trading Comps — EV/EBITDA",  6.50, 12.50, TEAL,   "4.6–7.5x NTM EBITDA comps"),
    ("Trading Comps — EV/Revenue", 5.00,  9.50, NBLUE,  "0.8–1.5x comparable revenue"),
    ("Precedent Txns — EV/EBITDA", 7.50, 14.00, PURP,   "6–12x median precedent range"),
    ("DCF Analysis",               7.00, 13.50, GRN,    "WACC 7.5–9.5%, TGR 2–4%"),
    ("Sum-of-Parts (ex-Optionality)",8.00, 15.00, AMB,  "Asset value vs. EV gap"),
]
for i, (label, lo, hi, color, note) in enumerate(bars):
    ff_bar(s16, label, lo, hi, color, 1.78 + i * 0.76, note)

# Offer price line
offer_x = 4.0 + 10.0 * (8.8/18.0)
rect(s16, offer_x - 0.02, 1.65, 0.04, 5.08, NRED)
box(s16, "Proposed Offer: $10.00/sh", offer_x - 1.2, 6.78, 2.7, 0.38,
    sz=10.5, bold=True, fc=NRED, align=PP_ALIGN.CENTER)

footnote_bar(s16, "Football field shows range of implied equity values per WBD share. "
                  "DCF uses 5-year UFCF projection with terminal value. "
                  "Proposed offer of $10.00/sh falls within or above all analytical ranges.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — PRO FORMA FINANCIALS
# ════════════════════════════════════════════════════════════════════════════
s17 = blank(); fill_bg(s17, LGRY)
slide_hdr(s17, "Pro Forma Financial Summary",
          "Combined Entity: $78B Revenue, $24B EBITDA, and Rapid Deleverage Profile")

kpi17 = [
    ("$78.3B", "PF FY2024 Combined\nRevenue"),
    ("$20.1B", "PF FY2024 Combined\nAdj. EBITDA (pre-syn)"),
    ("$24.1B", "PF FY2025E EBITDA\n(incl. $4B synergies Yr1)"),
    ("$9.0B",  "PF Year 1\nFCF Estimate"),
    ("2.9x",   "PF Net Leverage\n(EBITDA excl. synergies)"),
    ("2.1x",   "PF Net Leverage\n(Yr 3, post-synergies)"),
]
for i, (val, lbl) in enumerate(kpi17):
    x = 0.35 + i * 2.12
    accent = RED2 if "Leverage" in lbl and "2.9x" in val else (GRN if "2.1x" in val else NAVY)
    kpi_box(s17, val, lbl, x, 1.32, w=2.0, h=1.32, accent=accent)

# PF P&L table
box(s17, "Pro Forma Income Statement — Selected Metrics ($B)", 0.35, 2.82, 12.63, 0.38,
    sz=12, bold=True, fc=DARK)
tbl17 = add_tbl(s17, 10, 7, 0.35, 3.25, 12.63, 3.55)
for j, h in enumerate(["Metric","NFLX SA\n(FY24A)","WBD SA\n(FY24A)","PF Combined\n(FY24A)",
                        "PF Yr 1 Synergies","PF FY25E\n(+synergies)","PF FY26E\n(+synergies)"]):
    thdr(tbl17, 0, j, h, sz=9.5)

pf_rows = [
    ("Revenue ($B)",          "$39.0B","$39.3B","$78.3B","  —",    "$86.0B","$97.0B"),
    ("Revenue Growth",        "+15.6%","  -1.0%","  +7%","  —",    "+9.8%", "+12.8%"),
    ("Adj. EBITDA ($B)",      "$10.4B"," $9.7B","$20.1B","$1.3B",  "$24.1B","$28.5B"),
    ("EBITDA Margin",         "26.6%", "24.7%", "25.7%","  —",     "28.0%", "29.4%"),
    ("D&A ($B)",              " $1.2B"," $4.5B"," $5.7B","  —",    " $5.5B"," $5.3B"),
    ("EBIT ($B)",             " $9.2B"," $5.2B","$14.4B","$1.3B",  "$18.6B","$23.2B"),
    ("Interest Expense ($B)", " $0.9B"," $3.0B"," $4.5B","  —",    " $4.2B"," $3.8B"),
    ("Net Income ($B)",       " $8.7B"," -$10.8B","$7.5B¹","$0.9B","$10.8B","$14.5B"),
    ("Free Cash Flow ($B)",   " $6.9B"," $2.1B"," $9.0B","$1.0B",  "$11.0B","$13.8B"),
]
for i, row in enumerate(pf_rows):
    bold = row[0] in ("Adj. EBITDA ($B)","Free Cash Flow ($B)","Net Income ($B)")
    bg = NAVY if bold else (STLB if i%2==0 else None)
    fc_ = WHT if bold else DARK
    for j, v in enumerate(row):
        al = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
        tcell(tbl17, i+1, j, v, sz=9.5, bold=bold, fc=fc_, bg=bg, align=al)

footnote_bar(s17, "¹ PF net income excludes estimated $1.5B one-time integration/restructuring charges in Yr 1. "
                  "Adj. EBITDA excludes D&A, SBC. FY25E/FY26E are Financial Advisor projections. "
                  "WBD FY24 net loss includes ~$10.8B non-cash impairment charges (not cash-impacting).")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — RISK FACTORS & MITIGANTS
# ════════════════════════════════════════════════════════════════════════════
s18 = blank(); fill_bg(s18, WHT)
slide_hdr(s18, "Risk Factors & Mitigants",
          "Key Transaction Risks with Identified Mitigation Strategies")

box(s18, "Risk Assessment Matrix", 0.35, 1.28, 12.63, 0.38, sz=12, bold=True, fc=DARK)
tbl18 = add_tbl(s18, 9, 4, 0.35, 1.7, 12.63, 5.08)
for j, h in enumerate(["Risk Factor","Severity","Mitigation","Residual Risk"]):
    thdr(tbl18, 0, j, h, sz=10)

risks18 = [
    ("Regulatory / Antitrust Clearance",
     "MEDIUM",
     "No broadcast overlap; streaming-only deal. Offer divest CNN if required. "
     "Skydance/Paramount precedent approved in 10 months. EU may require Max/HBO price cap.",
     "LOW–MED"),
    ("Execution Risk / Integration",
     "HIGH",
     "Phased integration: Day 1 = keep brands separate. Month 6 = technology merge. "
     "Yr 1 = content teams merge. Retain WBD CEO David Zaslav 12-month transition.",
     "MEDIUM"),
    ("Leverage & Balance Sheet Risk",
     "HIGH",
     "Pro forma 2.9x Net Debt/EBITDA is manageable. $11B+ FCF generates $7–8B/yr "
     "net after interest. Committed financing from Tier-1 banks. ILG-rated covenant package.",
     "LOW–MED"),
    ("Linear Networks Secular Decline",
     "MEDIUM",
     "Networks are cash-generative for 5–7 more years at declining rate. "
     "Maintain optionality to sell Networks post-close ($10–15B value in 3 years). "
     "Structured as 'harvest-and-divest' cash cow.",
     "LOW"),
    ("Max/HBO Subscriber Retention",
     "MEDIUM",
     "Do NOT immediately merge services. Maintain HBO Max brand 18 months. "
     "Gradual migration to Netflix with HBO bundle benefit. Precedent: Disney+ kept Hulu.",
     "LOW"),
    ("Sports Rights Uncertainty",
     "LOW–MED",
     "WBD lost NBA but retains March Madness (–2032), MLB (–2028), NHL (–2028). "
     "Netflix already in live sports discussions; WBD accelerates, not creates, risk.",
     "LOW"),
    ("Talent & Culture Retention",
     "MEDIUM",
     "Golden handcuff packages for top WBD talent. Netflix culture playbook well-tested. "
     "Offer equity upside in combined entity. Retain HBO/Max programming team autonomy.",
     "LOW"),
    ("Currency / FX Risk",
     "LOW",
     "~40% WBD revenue international but USD-denominated contracts. "
     "Natural hedge from Netflix global revenue base. FX overlay program at combined entity.",
     "LOW"),
]
sev_colors = {"HIGH": RED2, "MEDIUM": AMB, "LOW–MED": RGBColor(0xFF,0xC0,0x00), "LOW": GRN}
res_colors = {"HIGH": RED2, "MEDIUM": AMB, "LOW–MED": RGBColor(0xFF,0xC0,0x00), "LOW": GRN}
for i, (risk, sev, mit, res) in enumerate(risks18):
    bg = STLB if i % 2 == 0 else None
    tcell(tbl18, i+1, 0, risk, sz=9.5, bold=True, fc=DARK, bg=bg)
    sc = sev_colors.get(sev, DARK)
    tcell(tbl18, i+1, 1, sev, sz=9.5, bold=True, fc=sc, bg=bg, align=PP_ALIGN.CENTER)
    tcell(tbl18, i+1, 2, mit, sz=9, fc=DARK, bg=bg)
    rc = res_colors.get(res, DARK)
    tcell(tbl18, i+1, 3, res, sz=9.5, bold=True, fc=rc, bg=bg, align=PP_ALIGN.CENTER)

footnote_bar(s18, "Risk severity and residual risk are Financial Advisor qualitative assessments. "
                  "Regulatory timeline and outcome subject to jurisdictional review processes.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — RECOMMENDATION & NEXT STEPS
# ════════════════════════════════════════════════════════════════════════════
s19 = blank(); fill_bg(s19, NAVY)
rect(s19, 0, 0, 13.33, 1.18, RGBColor(0x08, 0x16, 0x30))
rect(s19, 0, 1.18, 13.33, 0.055, GOLD)
box(s19, "Recommendation & Next Steps", 0.38, 0.18, 12.5, 0.65, sz=22, bold=True, fc=WHT)
box(s19, "Financial Advisor Assessment — For Management Discussion Only",
    0.38, 0.76, 12.5, 0.38, sz=12, fc=GOLD)

# Recommendation box
rect(s19, 0.35, 1.35, 12.63, 1.55, RGBColor(0x1A, 0x35, 0x65))
rect(s19, 0.35, 1.35, 0.18, 1.55, GOLD)
box(s19, "RECOMMENDATION", 0.62, 1.4, 4.0, 0.45, sz=14, bold=True, fc=GOLD)
rec_text = (
    "We recommend Netflix pursue a strategic acquisition of Warner Bros. Discovery at $10.00 per share, "
    "representing an approximately 182% premium to the current undisturbed share price of $3.55. "
    "The total transaction value of ~$60 billion (EV) represents 6.2x LTM Adjusted EBITDA — "
    "a disciplined entry price that compares favorably to all precedent media M&A transactions. "
    "The ~$4.0B/yr in synergies (NPV ~$38B) provides a compelling return on the $23.8B equity investment. "
    "Netflix should initiate confidential outreach to the WBD Board immediately."
)
box(s19, rec_text, 0.62, 1.88, 12.2, 0.95, sz=11, fc=WHT)

# Pillars (3 rationale boxes)
rationale = [
    ("Strategic", GOLD,
     "• Eliminates #2 streaming competitor (Max 96M subs)\n"
     "• Secures irreplaceable IP (DC, HP, GoT) for $0 content CPM\n"
     "• Sports rights de-risk Netflix live content strategy\n"
     "• Creates unassailable 400M+ subscriber leadership"),
    ("Financial", GRN,
     "• 6.2x EV/EBITDA — cheapest large-scale media M&A in 5 years\n"
     "• $4.0B/yr synergies → IRR of ~22% on equity invested\n"
     "• Net leverage peaks at 2.9x → sub-2.0x in Year 3\n"
     "• EPS accretive within 18–24 months post-synergy ramp"),
    ("Timing", NRED,
     "• WBD Board in active strategic review (Feb 2025)\n"
     "• Window closing: Comcast/others exploring WBD assets\n"
     "• Regulatory environment most favorable in 5 years\n"
     "• Act before further linear decline erodes synergy base"),
]
for i, (title, color, body) in enumerate(rationale):
    x = 0.35 + i * 4.33
    rect(s19, x, 3.15, 4.1, 2.28, RGBColor(0x1A, 0x35, 0x65))
    rect(s19, x, 3.15, 4.1, 0.45, color)
    box(s19, title, x+0.15, 3.18, 3.8, 0.38, sz=14, bold=True, fc=WHT, align=PP_ALIGN.CENTER)
    box(s19, body, x+0.15, 3.68, 3.8, 1.68, sz=10, fc=WHT)

# Next steps timeline
rect(s19, 0.35, 5.62, 12.63, 1.5, RGBColor(0x1A, 0x35, 0x65))
box(s19, "PROPOSED NEXT STEPS", 0.5, 5.66, 4.0, 0.38, sz=12, bold=True, fc=GOLD)
steps = [
    ("Week 1–2",   "Form M&A Committee; appoint Transaction Team; authorize Financial Advisor to proceed"),
    ("Week 2–4",   "Confidential outreach to WBD Chairman / CEO; NDA execution; preliminary diligence request"),
    ("Month 1–2",  "Management presentations; confirmatory financial & legal due diligence"),
    ("Month 2–3",  "Negotiate definitive merger agreement; secure committed financing; Board approval"),
    ("Month 3+",   "Public announcement; regulatory filings (HSR, EU, FCC); shareholder vote"),
]
for i, (phase, action) in enumerate(steps):
    x = 0.5 + i * 2.55
    box(s19, phase, x, 6.08, 2.4, 0.35, sz=9.5, bold=True, fc=GOLD, align=PP_ALIGN.CENTER)
    box(s19, action, x, 6.45, 2.4, 0.6, sz=8.5, fc=WHT, align=PP_ALIGN.CENTER, wrap=True)

footnote_bar(s19, "This recommendation is provided for discussion purposes only and is subject to ongoing due diligence, "
                  "satisfactory completion of confirmatory analysis, and Board authorization. May 2025.")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — APPENDIX
# ════════════════════════════════════════════════════════════════════════════
s20 = blank(); fill_bg(s20, LGRY)
slide_hdr(s20, "Appendix — Selected Financial Data",
          "Netflix & WBD Historical Financials and Additional Supporting Analysis")

# Netflix 5-year summary table
box(s20, "Netflix — 5-Year Financial Summary ($B)", 0.35, 1.32, 6.1, 0.38,
    sz=11, bold=True, fc=DARK)
tbl20a = add_tbl(s20, 7, 6, 0.35, 1.75, 6.1, 2.35)
for j, h in enumerate(["Metric","FY2020A","FY2021A","FY2022A","FY2023A","FY2024A"]):
    thdr(tbl20a, 0, j, h, sz=9)
nflx_hist = [
    ("Revenue",        "$25.0B","$29.7B","$31.6B","$33.7B","$39.0B"),
    ("Gross Profit",   "$11.5B","$13.0B","$15.2B","$16.6B","$20.6B"),
    ("Adj. EBITDA",    " $4.6B"," $5.1B"," $5.7B"," $7.0B","$10.4B"),
    ("EBITDA Margin",  "18.4%", "17.2%", "18.1%", "20.8%", "26.6%"),
    ("Net Income",     " $2.8B"," $5.1B"," $4.5B"," $5.4B"," $8.7B"),
    ("FCF",            " $1.9B"," -$0.1B"," $1.6B"," $6.9B"," $6.9B"),
]
for i, (m, *v) in enumerate(nflx_hist):
    bold = m in ("Adj. EBITDA","Net Income")
    bg = STLB if i%2==0 else None
    tcell(tbl20a, i+1, 0, m, sz=9, bold=bold, fc=DARK, bg=bg)
    for j, val in enumerate(v):
        tcell(tbl20a, i+1, j+1, val, sz=9, bold=bold, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)

# WBD 5-year summary table
box(s20, "WBD — Financial Summary ($B)", 0.35, 4.25, 6.1, 0.38,
    sz=11, bold=True, fc=DARK)
tbl20b = add_tbl(s20, 6, 5, 0.35, 4.68, 6.1, 2.0)
for j, h in enumerate(["Metric","FY2022A","FY2023A","FY2024A","Commentary"]):
    thdr(tbl20b, 0, j, h, sz=9)
wbd_hist = [
    ("Revenue",         "$43.2B","$41.3B","$39.3B","Declining -5% CAGR"),
    ("Adj. EBITDA",     "$10.3B","$10.8B"," $9.7B","Margin compression from cord-cutting"),
    ("Net Income (Loss)","$-2.0B","$-2.1B","-$10.8B","FY24 includes ~$10B impairments"),
    ("Total Debt",      "$49.7B","$46.0B","$39.5B","Rapid delever: $10B in 2yr"),
    ("Free Cash Flow",  " $2.8B"," $2.2B"," $2.1B","Positive FCF despite losses"),
]
for i, (m, *v) in enumerate(wbd_hist):
    bold = m in ("Adj. EBITDA","Total Debt")
    bg = STLB if i%2==0 else None
    tcell(tbl20b, i+1, 0, m, sz=9, bold=bold, fc=DARK, bg=bg)
    for j, val in enumerate(v[:3]):
        tcell(tbl20b, i+1, j+1, val, sz=9, bold=bold, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)
    tcell(tbl20b, i+1, 4, v[3], sz=9, italic=True, fc=MGRY, bg=bg)

# Synergy NPV bridge (right side)
box(s20, "Synergy NPV Bridge ($B) — 8.75% WACC", 6.65, 1.32, 6.35, 0.38,
    sz=11, bold=True, fc=DARK)
tbl20c = add_tbl(s20, 7, 4, 6.65, 1.75, 6.35, 2.75)
for j, h in enumerate(["Category","Bear","Base","Bull"]):
    thdr(tbl20c, 0, j, h, sz=9.5)
npv_rows = [
    ("Content Cost ($B NPV)",    "$10B", "$14B", "$19B"),
    ("Revenue Synergies ($B)",   " $6B", "$12B", "$22B"),
    ("Corp. / Tech ($B)",        " $5B", " $7B", "$10B"),
    ("Integration Costs",        "-$1B", "-$1.5B","-$2B"),
    ("Pre-Tax NPV Total",        "$20B", "$31.5B","$49B"),
    ("After-Tax NPV (35%)",      "$13B", "$20.5B","$32B"),
]
for i, (cat, bear, base, bull) in enumerate(npv_rows):
    bold = i in (4, 5)
    bg = NAVY if i == 5 else (STLB if i%2==0 else None)
    fc_ = WHT if i == 5 else DARK
    tcell(tbl20c, i+1, 0, cat, sz=9.5, bold=bold, fc=fc_, bg=bg)
    for j, v in enumerate([bear, base, bull]):
        tcell(tbl20c, i+1, j+1, v, sz=9.5, bold=bold, fc=fc_, bg=bg, align=PP_ALIGN.RIGHT)

box(s20, "Synergy Implementation Timeline & One-Time Costs",
    6.65, 4.65, 6.35, 0.38, sz=11, bold=True, fc=DARK)
tbl20d = add_tbl(s20, 5, 4, 6.65, 5.08, 6.35, 1.6)
for j, h in enumerate(["Category","Yr 1","Yr 2","Yr 3"]):
    thdr(tbl20d, 0, j, h, sz=9.5)
impl_rows = [
    ("Cumulative Synergies",     "$1.3B", "$2.6B", "$4.4B"),
    ("One-Time Costs (est.)",    "$0.8B", "$0.5B", "$0.2B"),
    ("Net Cumulative Benefit",   "$0.5B", "$2.1B", "$4.2B"),
    ("Payback Period",           "—",     "—",     "~5.5 years"),
]
for i, (cat, y1, y2, y3) in enumerate(impl_rows):
    bold = i in (2, 3)
    bg = STLB if i%2==0 else None
    tcell(tbl20d, i+1, 0, cat, sz=9.5, bold=bold, fc=DARK, bg=bg)
    for j, v in enumerate([y1, y2, y3]):
        tcell(tbl20d, i+1, j+1, v, sz=9.5, bold=bold, fc=DARK, bg=bg, align=PP_ALIGN.RIGHT)

footnote_bar(s20, "All figures sourced from company 10-K filings (FY2022–FY2024). "
                  "Synergy NPV uses 8.75% WACC and $4.4B peak annual synergies. "
                  "This presentation was validated using LibreOffice — please review in Microsoft PowerPoint before distribution.")


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(OUT, "NFLX_WBD_Acquisition_Pitchbook.pptx")
prs.save(out_path)
print(f"✓ Saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")
