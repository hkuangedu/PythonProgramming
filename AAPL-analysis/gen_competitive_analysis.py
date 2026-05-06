from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

APPLE_GRAY = RGBColor(0x1D, 0x1D, 0x1F)
APPLE_BLUE = RGBColor(0x00, 0x71, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF7)
MID_GRAY = RGBColor(0x86, 0x86, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height, font_size=18, bold=False,
                color=APPLE_GRAY, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s1 = blank_slide(prs)
fill_bg(s1, APPLE_GRAY)
add_rect(s1, 0, 0, 13.33, 7.5, APPLE_GRAY)
add_rect(s1, 0, 5.8, 13.33, 0.08, APPLE_BLUE)
add_textbox(s1, "APPLE (AAPL)", 1, 1.5, 11, 1.2, font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s1, "Competitive Landscape Analysis", 1, 2.9, 11, 0.8, font_size=26, color=RGBColor(0xA0,0xA0,0xA8), align=PP_ALIGN.CENTER)
add_textbox(s1, f"Prepared {datetime.date.today().strftime('%B %Y')}  |  Technology Sector  |  NASDAQ: AAPL",
            1, 4.2, 11, 0.5, font_size=14, color=MID_GRAY, align=PP_ALIGN.CENTER)

# ── Slide 2: Market Overview ────────────────────────────────────────────────
s2 = blank_slide(prs)
fill_bg(s2, LIGHT_GRAY)
add_rect(s2, 0, 0, 13.33, 1.1, APPLE_GRAY)
add_textbox(s2, "MARKET OVERVIEW & APPLE'S POSITION", 0.3, 0.25, 12, 0.7, font_size=22, bold=True, color=WHITE)

stats = [
    ("$3.3T", "Market Cap\n(#1 globally)"),
    ("$391B", "FY2024\nRevenue"),
    ("31%", "Net Profit\nMargin"),
    ("2.2B+", "Active\nDevices"),
    ("$110B+", "Annual\nR&D + CapEx"),
]
for i, (val, label) in enumerate(stats):
    x = 0.4 + i * 2.55
    add_rect(s2, x, 1.4, 2.2, 1.8, WHITE)
    add_textbox(s2, val, x, 1.55, 2.2, 0.8, font_size=28, bold=True, color=APPLE_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(s2, label, x, 2.35, 2.2, 0.7, font_size=12, color=MID_GRAY, align=PP_ALIGN.CENTER)

bullets = [
    "• Apple operates across five segments: iPhone (52%), Services (22%), Mac (8%), iPad (7%), Wearables/Home (11%)",
    "• Services segment (App Store, iCloud, Apple TV+, Apple Pay) growing at ~14% YoY — highest-margin business",
    "• Geographic mix: Americas 42%, Europe 25%, Greater China 19%, Rest of Asia Pacific 14%",
    "• Installed base of 2.2B+ active devices creates a durable, high-switching-cost ecosystem",
]
for i, b in enumerate(bullets):
    add_textbox(s2, b, 0.4, 3.45 + i * 0.7, 12.5, 0.6, font_size=13, color=APPLE_GRAY)

# ── Slide 3: Key Competitors ────────────────────────────────────────────────
s3 = blank_slide(prs)
fill_bg(s3, WHITE)
add_rect(s3, 0, 0, 13.33, 1.1, APPLE_GRAY)
add_textbox(s3, "KEY COMPETITORS", 0.3, 0.25, 12, 0.7, font_size=22, bold=True, color=WHITE)

competitors = [
    ("Microsoft\n(MSFT)", "$3.1T", "Cloud/AI/Enterprise", "Azure, Office 365, Copilot AI suite", APPLE_BLUE),
    ("Alphabet\n(GOOGL)", "$2.0T", "Search/Cloud/AI", "Android, Google Cloud, Gemini AI", RGBColor(0x34,0xA8,0x53)),
    ("Samsung", "$0.3T", "Hardware/Semicon.", "Galaxy devices, OLED displays, DRAM", RGBColor(0x03,0x78,0xFF)),
    ("Meta\n(META)", "$1.4T", "Social/AR/VR", "Vision Pro rival, Ray-Ban glasses", RGBColor(0x04,0x5D,0xE9)),
    ("Amazon\n(AMZN)", "$2.2T", "Cloud/Devices", "AWS, Alexa, Echo, Fire TV", RGBColor(0xFF,0x99,0x00)),
]
colors_row = [RGBColor(0xE8,0xF0,0xFE), RGBColor(0xE6,0xF4,0xEA), RGBColor(0xE3,0xF2,0xFF),
              RGBColor(0xE8,0xEA,0xFF), RGBColor(0xFF,0xF0,0xCC)]

headers = ["Company", "Mkt Cap", "Core Threat", "Key Products/Initiatives"]
col_w = [1.8, 1.2, 2.2, 5.8]
col_x = [0.3, 2.15, 3.4, 5.7]
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_textbox(s3, h, x, 1.2, w, 0.4, font_size=12, bold=True, color=MID_GRAY)

for i, (name, mktcap, threat, products, color) in enumerate(competitors):
    y = 1.65 + i * 1.0
    add_rect(s3, 0.3, y, 12.7, 0.88, LIGHT_GRAY if i % 2 == 0 else WHITE)
    data = [name, mktcap, threat, products]
    for j, (val, x, w) in enumerate(zip(data, col_x, col_w)):
        fs = 13 if j != 0 else 13
        add_textbox(s3, val, x, y + 0.05, w, 0.78, font_size=fs, bold=(j == 0), color=APPLE_GRAY)

# ── Slide 4: Competitive Moats ──────────────────────────────────────────────
s4 = blank_slide(prs)
fill_bg(s4, LIGHT_GRAY)
add_rect(s4, 0, 0, 13.33, 1.1, APPLE_GRAY)
add_textbox(s4, "APPLE'S COMPETITIVE MOATS", 0.3, 0.25, 12, 0.7, font_size=22, bold=True, color=WHITE)

moats = [
    ("Ecosystem Lock-In", APPLE_BLUE,
     "iMessage, AirDrop, Handoff, iCloud create seamless device interoperability. "
     "Switching cost estimated at $1,000–$1,500 per user. Churn rate <5%/year."),
    ("Brand & Premium Pricing", RGBColor(0x30,0xB0,0xC7),
     "Consistent #1 in brand loyalty (ACSI). ASP of $850+ vs. Android ~$300. "
     "Premium pricing power sustained across economic cycles."),
    ("Vertical Integration", RGBColor(0x34,0xC7,0x59),
     "Custom silicon (M-series, A-series chips) provides 2–3 year CPU/GPU performance lead. "
     "Owns chip design, software, services, and retail distribution."),
    ("Services Flywheel", RGBColor(0xFF,0x95,0x00),
     "$96B+ annual services revenue growing 14% YoY at ~75% gross margin. "
     "App Store, Apple Pay, iCloud, Apple TV+, AppleCare compound with installed base."),
    ("Supply Chain Mastery", RGBColor(0xFF,0x45,0x22),
     "Multi-decade supplier relationships, TSMC allocation priority, and $110B+ annual "
     "procurement give unmatched cost structure and component security."),
]
for i, (title, color, desc) in enumerate(moats):
    col = i % 2
    row = i // 2
    x = 0.3 + col * 6.5
    y = 1.25 + row * 2.2
    w = 6.1
    add_rect(s4, x, y, w, 1.95, WHITE)
    add_rect(s4, x, y, 0.18, 1.95, color)
    add_textbox(s4, title, x + 0.3, y + 0.1, w - 0.5, 0.45, font_size=14, bold=True, color=APPLE_GRAY)
    add_textbox(s4, desc, x + 0.3, y + 0.55, w - 0.5, 1.3, font_size=11, color=MID_GRAY)

# ── Slide 5: Comparative Metrics ────────────────────────────────────────────
s5 = blank_slide(prs)
fill_bg(s5, WHITE)
add_rect(s5, 0, 0, 13.33, 1.1, APPLE_GRAY)
add_textbox(s5, "COMPARATIVE ANALYSIS — KEY METRICS", 0.3, 0.25, 12, 0.7, font_size=22, bold=True, color=WHITE)

table_data = [
    ["Metric",         "Apple",  "Microsoft", "Alphabet", "Meta",   "Amazon"],
    ["Revenue (FY24)", "$391B",  "$245B",     "$350B",    "$165B",  "$620B"],
    ["Revenue Growth", "+6%",    "+16%",      "+14%",     "+22%",   "+11%"],
    ["Gross Margin",   "46%",    "70%",       "57%",      "81%",    "49%"],
    ["Net Margin",     "31%",    "35%",       "26%",      "38%",    "9%"],
    ["EV/EBITDA",      "27x",    "31x",       "18x",      "22x",    "24x"],
    ["P/E (NTM)",      "32x",    "36x",       "21x",      "25x",    "44x"],
    ["FCF Yield",      "3.1%",   "2.6%",      "4.2%",     "3.8%",   "2.1%"],
]
col_widths = [2.4, 1.7, 1.7, 1.7, 1.7, 1.7]
col_starts = [0.3]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

for i, row in enumerate(table_data):
    y = 1.2 + i * 0.68
    bg = APPLE_GRAY if i == 0 else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    add_rect(s5, 0.3, y, 12.6, 0.65, bg)
    for j, (cell, x, w) in enumerate(zip(row, col_starts, col_widths)):
        fc = WHITE if i == 0 else (APPLE_BLUE if j == 1 and i > 0 else APPLE_GRAY)
        add_textbox(s5, cell, x + 0.05, y + 0.08, w - 0.1, 0.5,
                    font_size=12, bold=(i == 0 or j == 0), color=fc)

add_textbox(s5, "Source: Bloomberg, FactSet consensus estimates. FY2024 actuals.",
            0.3, 6.8, 12, 0.4, font_size=10, color=MID_GRAY, italic=True)

# ── Slide 6: Strategic Risks & Opportunities ────────────────────────────────
s6 = blank_slide(prs)
fill_bg(s6, LIGHT_GRAY)
add_rect(s6, 0, 0, 13.33, 1.1, APPLE_GRAY)
add_textbox(s6, "STRATEGIC RISKS & OPPORTUNITIES", 0.3, 0.25, 12, 0.7, font_size=22, bold=True, color=WHITE)

risks = [
    ("China Revenue Concentration", "19% of revenue from Greater China; geopolitical risk and local brand competition from Huawei resurgence."),
    ("Regulatory / Antitrust", "EU Digital Markets Act, DOJ App Store investigation. Potential 30% fee reduction could impact $30B+ Services revenue."),
    ("AI Hardware Catch-Up", "On-device AI (Apple Intelligence) lags Qualcomm/Snapdragon X Elite in raw AI performance benchmarks."),
    ("Services Growth Dependency", "Multiple expansion requires Services to sustain 14%+ growth; any deceleration re-rates the stock materially."),
]
opps = [
    ("Apple Intelligence / GenAI", "On-device LLM integration across 2.2B devices — largest private AI distribution channel. ChatGPT partnership."),
    ("Financial Services Expansion", "Apple Card, Apple Pay Later, Savings Account — path to $50B+ fintech revenue by 2030."),
    ("Vision Pro Platform", "visionOS ecosystem still nascent; enterprise adoption in healthcare/AEC could open new TAM."),
    ("India Manufacturing & Market", "Shifting 25%+ production to India; domestic market of 1.4B with rising middle class."),
]

add_textbox(s6, "KEY RISKS", 0.3, 1.2, 6, 0.4, font_size=14, bold=True, color=RGBColor(0xD7,0x00,0x15))
add_textbox(s6, "KEY OPPORTUNITIES", 6.8, 1.2, 6, 0.4, font_size=14, bold=True, color=RGBColor(0x34,0xC7,0x59))

for i, (title, desc) in enumerate(risks):
    y = 1.7 + i * 1.3
    add_rect(s6, 0.3, y, 6.2, 1.18, WHITE)
    add_rect(s6, 0.3, y, 0.15, 1.18, RGBColor(0xD7,0x00,0x15))
    add_textbox(s6, title, 0.55, y + 0.05, 5.85, 0.4, font_size=12, bold=True, color=APPLE_GRAY)
    add_textbox(s6, desc, 0.55, y + 0.45, 5.85, 0.65, font_size=11, color=MID_GRAY)

for i, (title, desc) in enumerate(opps):
    y = 1.7 + i * 1.3
    add_rect(s6, 6.8, y, 6.2, 1.18, WHITE)
    add_rect(s6, 6.8, y, 0.15, 1.18, RGBColor(0x34,0xC7,0x59))
    add_textbox(s6, title, 7.05, y + 0.05, 5.85, 0.4, font_size=12, bold=True, color=APPLE_GRAY)
    add_textbox(s6, desc, 7.05, y + 0.45, 5.85, 0.65, font_size=11, color=MID_GRAY)

# ── Slide 7: Summary / Investment Thesis ────────────────────────────────────
s7 = blank_slide(prs)
fill_bg(s7, APPLE_GRAY)
add_rect(s7, 0, 0, 13.33, 1.1, APPLE_BLUE)
add_textbox(s7, "SUMMARY & INVESTMENT THESIS", 0.3, 0.2, 12, 0.75, font_size=22, bold=True, color=WHITE)

add_textbox(s7, "Bull Case", 0.4, 1.25, 3.8, 0.45, font_size=15, bold=True, color=RGBColor(0x34,0xC7,0x59))
bull_pts = [
    "Services flywheel compounds with 2.2B device base",
    "Apple Intelligence drives upgrade supercycle (FY25–26)",
    "India manufacturing de-risks China concentration",
    "Buyback yield (~3.5%) + dividend support total return",
    "Vision Pro creates next computing platform option value",
]
for i, pt in enumerate(bull_pts):
    add_textbox(s7, f"+ {pt}", 0.4, 1.75 + i * 0.72, 3.8, 0.65, font_size=12, color=WHITE)

add_textbox(s7, "Bear Case", 4.5, 1.25, 3.8, 0.45, font_size=15, bold=True, color=RGBColor(0xFF,0x45,0x22))
bear_pts = [
    "32x P/E leaves little room for execution misses",
    "China risk (Huawei + regulatory) threatens top-line",
    "EU antitrust reduces App Store fee economics",
    "Slowing iPhone ASP growth limits revenue upside",
    "AI compute intensity may strain margin profile",
]
for i, pt in enumerate(bear_pts):
    add_textbox(s7, f"- {pt}", 4.5, 1.75 + i * 0.72, 3.8, 0.65, font_size=12, color=WHITE)

add_textbox(s7, "Verdict", 8.6, 1.25, 4.3, 0.45, font_size=15, bold=True, color=RGBColor(0xFF,0xD6,0x00))
add_rect(s7, 8.6, 1.75, 4.3, 3.6, RGBColor(0x2C,0x2C,0x2E))
verdict = (
    "Apple is the world's most valuable company for a reason: "
    "its ecosystem moat, vertical integration, and Services flywheel "
    "are structurally durable.\n\n"
    "At current multiples (32x NTM P/E), the stock prices in "
    "continued Services outperformance and an AI-driven device cycle. "
    "The risk/reward is balanced — compelling for long-term holders, "
    "but entry timing matters.\n\n"
    "Target price range: $205–$240 (base case)\n"
    "Bull case: $265+ | Bear case: $175"
)
add_textbox(s7, verdict, 8.75, 1.85, 3.95, 3.35, font_size=11, color=WHITE)

# Save
out = "/root/PythonProgramming/AAPL-analysis/AAPL_Competitive_Analysis.pptx"
prs.save(out)
print(f"Saved: {out}")
